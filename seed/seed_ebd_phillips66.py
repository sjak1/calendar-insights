"""
EBD Seeding Script - Phillips 66 Cloud Security
Creates a filled-out EBD document for Phillips 66 focusing on Cloud Security.
Run: python seed/seed_ebd_phillips66.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path


# Phillips 66 Cloud Security focused data
DUMMY_DATA = {
    "company_name": "Phillips 66",
    "industry": "Oil & Gas / Energy",
    "address": "2331 CityWest Blvd, Houston, TX 77042",
    "website": "www.phillips66.com",
    
    "meeting_date": "March 15, 2026 | 9:00 AM - 4:00 PM CST",
    "meeting_location": "Oracle Executive Briefing Center - Austin",
    "engagement_type": "1:1 Executive Briefing",
    
    "ebd_contact": "Michael Torres, Sr. Account Executive, michael.torres@oracle.com, (713) 555-4567",
    
    "oracle_attendees": """David Chen, VP Cloud Security Solutions, david.chen@oracle.com
Rachel Martinez, Sr. Director Infrastructure Security, rachel.martinez@oracle.com
Kevin Park, Principal Security Architect, kevin.park@oracle.com""",
    
    "customer_attendees": """William Anderson, CISO
Jennifer Hughes, VP IT Infrastructure
Robert Thompson, Director Cloud Operations
Sarah Mitchell, Head of Compliance & Risk
Marcus Lee, Senior Security Engineer""",
    
    "company_description": """Phillips 66 is a $150B diversified energy manufacturing and logistics company. Headquartered in Houston, Texas, they operate 13 refineries, 22 terminals, and over 7,000 retail locations. With 14,000 employees across North America and Europe, Phillips 66 processes 2.2 million barrels of crude oil per day. They are in the midst of a major cloud migration initiative to modernize their IT infrastructure while maintaining strict security and compliance requirements for critical energy infrastructure.""",
    
    "meeting_objectives": """1. Evaluate Oracle Cloud Infrastructure (OCI) security capabilities for protecting critical energy infrastructure and operational technology (OT) systems
2. Understand Oracle's Zero Trust security architecture and how it applies to hybrid cloud environments
3. Review Oracle's compliance certifications (SOC 2, ISO 27001, NERC CIP) for energy sector requirements
4. Discuss Oracle's threat detection and response capabilities using AI/ML
5. Explore Oracle's data sovereignty and regional data residency options for international operations""",
    
    "business_challenges": """1. HYBRID CLOUD SECURITY: Managing security across on-premises SCADA systems, private cloud, and public cloud environments; current fragmented security tools cost $8M annually and create visibility gaps

2. RANSOMWARE THREAT: Energy sector is #1 target for ransomware attacks; recent industry peer suffered $45M attack; need advanced threat detection and immutable backups

3. COMPLIANCE COMPLEXITY: Must maintain NERC CIP, TSA Pipeline Security, and SOX compliance across 13 refineries; current manual audit processes take 6+ months and cost $5M annually

4. IDENTITY & ACCESS MANAGEMENT: 14,000 employees plus 8,000 contractors need secure access; legacy IAM system lacks modern MFA and privileged access management; recent audit found 2,400 orphaned accounts

5. OT/IT CONVERGENCE: Industrial control systems increasingly connected to IT networks; need secure architecture that protects operational technology without impacting refinery operations""",
    
    "it_strategy": """CEO Vision: "Secure digital transformation to become the safest, most efficient energy company by 2030"

IT Security Strategy:
- Zero Trust architecture adoption across all environments
- Cloud-first for new workloads with security-by-design
- Unified security operations center (SOC) with AI-driven threat detection
- Automated compliance monitoring and reporting
- Identity-centric security model with continuous verification
- 24/7 threat monitoring with <15 minute response SLA""",
    
    "customer_lifecycle": """SALES STAGE: Evaluation / Active Opportunity ($12.5M TCV)
- Security assessment completed Q4 2025
- RFP issued to Oracle, Microsoft Azure, and AWS
- Security deep-dive scheduled for March 2026
- Decision expected Q2 2026

WHY CONSIDERING ORACLE:
- Existing Oracle Database footprint (80% of enterprise databases)
- Oracle's integrated security stack appeal vs. point solutions
- Energy sector references (Chevron, Shell, BP)
- Oracle's government and defense security credentials

RECENT IMPLEMENTATIONS:
- Oracle Identity Cloud Service POC completed (positive results)
- Oracle Autonomous Database deployed for financial systems (2024)""",
    
    "account_status": """YELLOW - Opportunity at risk

POSITIVES:
- CISO (William Anderson) impressed with Oracle's security-first architecture
- VP IT Infrastructure prefers Oracle's integrated approach over Azure's fragmented tools
- Strong existing Oracle Database relationship

POTENTIAL DERAILERS:
- CFO concerned about Oracle's pricing vs. AWS/Azure
- Microsoft offering significant Azure credits ($5M)
- Board questioning why not standardize on Microsoft given O365 deployment
- Concerns about Oracle's energy sector security references""",
    
    "oracle_talking_points": """1. INTEGRATED SECURITY STACK: Demonstrate Oracle's unified security platform (OCI Security, Identity, Key Management) vs. competitors' fragmented approach - highlight $3M+ annual savings from tool consolidation at similar energy companies

2. ENERGY SECTOR EXPERTISE: Showcase Oracle's energy customers (Chevron achieved 60% faster threat response, Shell reduced compliance costs by 40%) and our NERC CIP / TSA compliance capabilities

3. AI-POWERED THREAT DETECTION: Present Oracle's Cloud Guard and Security Zones with autonomous threat detection - differentiate from AWS/Azure with embedded AI that requires no additional configuration""",
    
    # Attendee profiles
    "attendee_1_name": "William Anderson, CISO, william.anderson@phillips66.com",
    "attendee_1_perspective": "Primary decision maker for security investments. Laser-focused on protecting critical infrastructure from nation-state threats. Previously led security at Marathon Oil. Skeptical of cloud security but open to evidence. Wants to see real-world energy sector deployments.",
    "attendee_1_bio": "20+ years in cybersecurity, former NSA analyst. CISSP, CISM certified. Board member at Oil & Gas ISAC (Information Sharing and Analysis Center).",
    
    "attendee_2_name": "Jennifer Hughes, VP IT Infrastructure, jennifer.hughes@phillips66.com",
    "attendee_2_perspective": "Responsible for cloud migration strategy. Strong technical background. Concerned about operational continuity during migration. Prefers Oracle's integrated stack but facing pressure to consider Azure due to Microsoft relationship.",
    "attendee_2_bio": "Former Director at ExxonMobil IT. 18 years in energy sector infrastructure. MS in Computer Science from Texas A&M.",
    
    "attendee_3_name": "Robert Thompson, Director Cloud Operations, robert.thompson@phillips66.com",
    "attendee_3_perspective": "Hands-on technical leader managing current cloud footprint. Frustrated with current multi-cloud complexity. Wants unified security visibility. Will ask detailed technical questions about OCI networking and security controls.",
    "attendee_3_bio": "15 years in cloud infrastructure. Former AWS solutions architect. Kubernetes and cloud security certifications.",
    
    # Customer references
    "reference_1_name": "Chevron",
    "reference_1_industry": "Oil & Gas",
    "reference_1_products": "Oracle Cloud Infrastructure, Oracle Identity Cloud, Oracle Security Operations",
    "reference_1_summary": "Migrated 500+ critical applications to OCI with zero security incidents. Achieved 60% faster threat detection and response time. Reduced security tool sprawl from 45 to 12 solutions, saving $4.5M annually.",
    
    "reference_2_name": "Shell",
    "reference_2_industry": "Energy",
    "reference_2_products": "Oracle Cloud Infrastructure, Oracle Cloud Guard, Oracle Data Safe",
    "reference_2_summary": "Implemented Oracle's security stack across hybrid environment. Achieved NERC CIP compliance in 3 months vs. 12 months with previous solution. Automated 85% of compliance reporting, saving 10,000 person-hours annually.",
    
    "reference_3_name": "BP",
    "reference_3_industry": "Oil & Gas",
    "reference_3_products": "Oracle Cloud Infrastructure, Oracle Key Management, Oracle Identity Cloud",
    "reference_3_summary": "Deployed Oracle's Zero Trust architecture across 70 global locations. Reduced unauthorized access attempts by 95%. Achieved SOC 2 Type II and ISO 27001 certification with Oracle's compliance automation tools.",
    
    # Previous engagements
    "previous_engagement_1": "1:1 Executive Briefing | December 10, 2025 / Virtual | David Chen | William Anderson, Jennifer Hughes",
    "previous_engagement_2": "Oracle CloudWorld | September 2025 / Las Vegas | Mark Hurd Memorial Security Summit | William Anderson, Robert Thompson",
    
    # Account team
    "account_manager": "Michael Torres, Sr. Account Executive, michael.torres@oracle.com, (713) 555-4567",
    "other_account_manager": "Lisa Wong, Strategic Account Manager, lisa.wong@oracle.com, (713) 555-8901",
    "csm": "James Rodriguez, Customer Success Manager, james.rodriguez@oracle.com, (713) 555-2345",
    "executive_sponsor": "Sarah Mitchell, SVP Energy Sector Sales, sarah.mitchell@oracle.com",
    "partner": "Deloitte - Energy & Resources Cybersecurity Practice (Primary), Accenture Security (Secondary)",
    
    # Financials
    "deal_size_saas": "$5.2M",
    "deal_prob_saas": "55%",
    "deal_size_paas": "$7.3M", 
    "deal_prob_paas": "60%",
    "total_deal_size": "$12.5M",
    "tam": "$85M",
    "annual_spend": "$4.8M",
    "company_revenue": "$150B",
    "share_of_wallet": "3.2%",
}


def fill_table_cell(table, row, col, text):
    """Fill a table cell with text, preserving formatting."""
    cell = table.cell(row, col)
    if cell.text_frame.paragraphs:
        cell.text_frame.paragraphs[0].text = text


def fill_ebd_template(template_path: str, output_path: str, data: dict):
    """
    Fill the EBD template with provided data.
    """
    prs = Presentation(template_path)
    
    for slide_num, slide in enumerate(prs.slides, 1):
        # Skip first slide (instructions)
        if slide_num == 1:
            continue
            
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                
                # Slide 2: Main company info table
                if slide_num == 2:
                    # Fill company info
                    for row_idx, row in enumerate(table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            text = cell.text.strip().lower()
                            
                            # Company name and industry (first row)
                            if "company name" in text or (row_idx == 0 and col_idx == 0):
                                fill_table_cell(table, row_idx, col_idx, f"{data['company_name']} | {data['industry']}")
                            
                            # Address
                            if "address" in text or "terry avenue" in text.lower():
                                fill_table_cell(table, row_idx, col_idx, data['address'])
                            
                            # Website
                            if "website" in text or "www." in text:
                                fill_table_cell(table, row_idx, col_idx, data['website'])
                            
                            # Meeting date/time
                            if "meeting date" in text:
                                next_col = min(col_idx + 1, len(row.cells) - 1)
                                fill_table_cell(table, row_idx, next_col, data['meeting_date'])
                            
                            # Meeting location
                            if "meeting location" in text:
                                next_col = min(col_idx + 1, len(row.cells) - 1)
                                fill_table_cell(table, row_idx, next_col, data['meeting_location'])
                            
                            # Engagement type
                            if "engagement type" in text:
                                next_col = min(col_idx + 1, len(row.cells) - 1)
                                fill_table_cell(table, row_idx, next_col, data['engagement_type'])
                            
                            # EBD contact
                            if "ebd" in text and "contact" in text:
                                next_col = min(col_idx + 1, len(row.cells) - 1)
                                fill_table_cell(table, row_idx, next_col, data['ebd_contact'])
                            
                            # Oracle attendees
                            if "oracle meeting attendees" in text or "oracle attendees" in text:
                                next_col = min(col_idx + 1, len(row.cells) - 1)
                                fill_table_cell(table, row_idx, next_col, data['oracle_attendees'])
                            
                            # Customer attendees
                            if "customer meeting attendees" in text or "customer attendees" in text:
                                next_col = min(col_idx + 1, len(row.cells) - 1)
                                fill_table_cell(table, row_idx, next_col, data['customer_attendees'])
                            
                            # Company description
                            if "company description" in text:
                                next_col = min(col_idx + 1, len(row.cells) - 1)
                                fill_table_cell(table, row_idx, next_col, data['company_description'])
                            
                            # Meeting objectives
                            if "meeting objectives" in text or "customer expectations" in text:
                                next_col = min(col_idx + 1, len(row.cells) - 1)
                                fill_table_cell(table, row_idx, next_col, data['meeting_objectives'])
                            
                            # Business challenges
                            if "business challenges" in text or "where oracle can help" in text:
                                next_col = min(col_idx + 1, len(row.cells) - 1)
                                fill_table_cell(table, row_idx, next_col, data['business_challenges'])
                            
                            # IT Strategy
                            if "it strategy" in text or "business" in text and "strategy" in text:
                                next_col = min(col_idx + 1, len(row.cells) - 1)
                                fill_table_cell(table, row_idx, next_col, data['it_strategy'])
                            
                            # Customer lifecycle
                            if "customer lifecycle" in text or "sales process" in text:
                                next_col = min(col_idx + 1, len(row.cells) - 1)
                                fill_table_cell(table, row_idx, next_col, data['customer_lifecycle'])
                            
                            # Account status
                            if "account status" in text:
                                next_col = min(col_idx + 1, len(row.cells) - 1)
                                fill_table_cell(table, row_idx, next_col, data['account_status'])
                            
                            # Oracle talking points
                            if "oracle talking points" in text:
                                next_col = min(col_idx + 1, len(row.cells) - 1)
                                fill_table_cell(table, row_idx, next_col, data['oracle_talking_points'])
                
                # Slide 3: Account team and attendee profiles
                if slide_num == 3:
                    for row_idx, row in enumerate(table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            text = cell.text.strip().lower()
                            
                            # Previous engagements
                            if "1:1 executive briefing" in text and row_idx > 0:
                                fill_table_cell(table, row_idx, col_idx, data.get('previous_engagement_1', ''))
                            
                            # Account manager
                            if "scd" in text or "lam" in text:
                                next_col = min(col_idx + 1, len(row.cells) - 1)
                                fill_table_cell(table, row_idx, next_col, data['account_manager'])
                            
                            # Other account managers
                            if "other account manager" in text:
                                next_col = min(col_idx + 1, len(row.cells) - 1)
                                fill_table_cell(table, row_idx, next_col, data['other_account_manager'])
                            
                            # CSM
                            if "cse" in text or "csm" in text:
                                next_col = min(col_idx + 1, len(row.cells) - 1)
                                fill_table_cell(table, row_idx, next_col, data['csm'])
                            
                            # Executive sponsor
                            if "executive sponsor" in text:
                                next_col = min(col_idx + 1, len(row.cells) - 1)
                                fill_table_cell(table, row_idx, next_col, data['executive_sponsor'])
                            
                            # Partner
                            if "implementation partner" in text or "partner" in text:
                                next_col = min(col_idx + 1, len(row.cells) - 1)
                                fill_table_cell(table, row_idx, next_col, data['partner'])
                            
                            # Attendee profiles - look for patterns
                            if "ceo" in text or "cfo" in text or "cio" in text or "ciso" in text:
                                # This is an attendee row, fill with our data
                                if "ciso" in text or (row_idx < 3 and "attendee_1" in str(data)):
                                    fill_table_cell(table, row_idx, col_idx, data['attendee_1_name'])
                                    if col_idx + 1 < len(row.cells):
                                        fill_table_cell(table, row_idx, col_idx + 1, data['attendee_1_perspective'])
                
                # Slide 4: Customer references
                if slide_num == 4:
                    for row_idx, row in enumerate(table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            text = cell.text.strip().lower()
                            
                            if "customer name" in text:
                                next_col = min(col_idx + 1, len(row.cells) - 1)
                                # Determine which reference based on table position
                                ref_num = 1  # Default
                                fill_table_cell(table, row_idx, next_col, data.get(f'reference_{ref_num}_name', ''))
                            
                            if "customer industry" in text:
                                next_col = min(col_idx + 1, len(row.cells) - 1)
                                fill_table_cell(table, row_idx, next_col, data.get('reference_1_industry', ''))
                            
                            if "oracle product" in text:
                                next_col = min(col_idx + 1, len(row.cells) - 1)
                                fill_table_cell(table, row_idx, next_col, data.get('reference_1_products', ''))
                            
                            if "summary" in text:
                                next_col = min(col_idx + 1, len(row.cells) - 1)
                                fill_table_cell(table, row_idx, next_col, data.get('reference_1_summary', ''))
    
    prs.save(output_path)
    print(f"Created: {output_path}")


def main():
    # Use the existing template
    template_path = Path(__file__).parent.parent / "documents" / "ebd" / "EBD_template.pptx"
    
    if not template_path.exists():
        # Use Amazon as template if no template exists
        template_path = Path(__file__).parent.parent / "documents" / "ebd" / "EBD_Amazon_FILLED.pptx"
    
    if not template_path.exists():
        print(f"Error: No template found at {template_path}")
        return
    
    output_path = Path(__file__).parent.parent / "documents" / "ebd" / "EBD_Phillips66_FILLED.pptx"
    
    print(f"Using template: {template_path}")
    print(f"Creating Phillips 66 EBD with Cloud Security focus...")
    
    # Create the filled EBD
    fill_ebd_template(str(template_path), str(output_path), DUMMY_DATA)
    
    print(f"\nDone! Created: {output_path}")
    print("\nTo test extraction:")
    print(f"  python tools/extract_ebd.py {output_path}")


if __name__ == "__main__":
    main()

