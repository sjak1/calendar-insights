"""
Create a professional PowerPoint presentation on:
Difference between AI Chatbots, LLMs, and Autonomous AI Agents
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme - Corporate modern
DARK_BLUE = RGBColor(0x1a, 0x1a, 0x2e)
ACCENT_BLUE = RGBColor(0x41, 0x6d, 0x9c)
LIGHT_BLUE = RGBColor(0x5d, 0xa3, 0xd9)
ORANGE = RGBColor(0xe9, 0x7a, 0x3b)
GREEN = RGBColor(0x2e, 0xcc, 0x71)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xf0, 0xf0, 0xf0)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)

def add_title_slide(prs, title, subtitle):
    """Add a title slide with gradient-like background"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()
    
    # Accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.2), prs.slide_width, Inches(0.1))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items, accent_color=ACCENT_BLUE):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # Top accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(22)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(14)
        p.level = 0
    
    return slide

def add_comparison_slide(prs, title, left_title, left_items, right_title, right_items, left_color, right_color):
    """Add a two-column comparison slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Left column header
    left_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.6))
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = left_color
    left_header.line.fill.background()
    
    lh_tf = left_header.text_frame
    lh_tf.paragraphs[0].text = left_title
    lh_tf.paragraphs[0].font.size = Pt(20)
    lh_tf.paragraphs[0].font.bold = True
    lh_tf.paragraphs[0].font.color.rgb = WHITE
    lh_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    lh_tf.word_wrap = True
    
    # Left content
    left_box = slide.shapes.add_textbox(Inches(0.6), Inches(2), Inches(5.6), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(10)
    
    # Right column header
    right_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.2), Inches(5.8), Inches(0.6))
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = right_color
    right_header.line.fill.background()
    
    rh_tf = right_header.text_frame
    rh_tf.paragraphs[0].text = right_title
    rh_tf.paragraphs[0].font.size = Pt(20)
    rh_tf.paragraphs[0].font.bold = True
    rh_tf.paragraphs[0].font.color.rgb = WHITE
    rh_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    rh_tf.word_wrap = True
    
    # Right content
    right_box = slide.shapes.add_textbox(Inches(7.1), Inches(2), Inches(5.6), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(10)
    
    return slide

def add_evolution_slide(prs, title, stages):
    """Add an evolution/timeline slide showing progression"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Timeline line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.5), Inches(11.7), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()
    
    # Stage cards
    colors = [RGBColor(0x95, 0xa5, 0xa6), ACCENT_BLUE, GREEN]
    x_positions = [0.8, 5.0, 9.2]
    
    for i, (stage_title, stage_desc, icon) in enumerate(stages):
        # Circle node
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_positions[i] + 1.3), Inches(3.35), Inches(0.35), Inches(0.35))
        circle.fill.solid()
        circle.fill.fore_color.rgb = colors[i]
        circle.line.fill.background()
        
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_positions[i]), Inches(4.0), Inches(3.5), Inches(2.8))
        card.fill.solid()
        card.fill.fore_color.rgb = colors[i]
        card.line.fill.background()
        
        # Icon/emoji above title
        icon_box = slide.shapes.add_textbox(Inches(x_positions[i]), Inches(1.5), Inches(3.5), Inches(0.6))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(40)
        p.alignment = PP_ALIGN.CENTER
        
        # Stage title above card
        st_box = slide.shapes.add_textbox(Inches(x_positions[i]), Inches(2.2), Inches(3.5), Inches(0.6))
        tf = st_box.text_frame
        p = tf.paragraphs[0]
        p.text = stage_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = colors[i]
        p.alignment = PP_ALIGN.CENTER
        
        # Description inside card
        desc_box = slide.shapes.add_textbox(Inches(x_positions[i] + 0.15), Inches(4.15), Inches(3.2), Inches(2.5))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = stage_desc
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # Arrows between stages
    for i in range(2):
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x_positions[i] + 3.6), Inches(3.3), Inches(1.0), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ORANGE
        arrow.line.fill.background()
    
    return slide

def add_three_column_slide(prs, title, col1, col2, col3):
    """Add a three-column comparison slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    columns = [col1, col2, col3]
    colors = [RGBColor(0x95, 0xa5, 0xa6), ACCENT_BLUE, GREEN]
    x_positions = [0.4, 4.6, 8.8]
    
    for i, (col_title, col_items) in enumerate(columns):
        # Header
        header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_positions[i]), Inches(0.9), Inches(4), Inches(0.55))
        header.fill.solid()
        header.fill.fore_color.rgb = colors[i]
        header.line.fill.background()
        
        h_tf = header.text_frame
        h_tf.paragraphs[0].text = col_title
        h_tf.paragraphs[0].font.size = Pt(18)
        h_tf.paragraphs[0].font.bold = True
        h_tf.paragraphs[0].font.color.rgb = WHITE
        h_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(x_positions[i] + 0.1), Inches(1.6), Inches(3.8), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for j, item in enumerate(col_items):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(15)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(8)
    
    return slide


def add_limitation_solution_slide(prs, title, limitation, solution_title, solution_items, limit_color, solution_color):
    """Slide showing limitation on left, solution on right"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Limitation box (left)
    limit_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(5.5), Inches(5.8))
    limit_box.fill.solid()
    limit_box.fill.fore_color.rgb = RGBColor(0xff, 0xeb, 0xeb)  # Light red
    limit_box.line.color.rgb = RGBColor(0xe7, 0x4c, 0x3c)
    limit_box.line.width = Pt(2)
    
    # Limitation header
    lh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.4), Inches(5.1), Inches(0.5))
    lh.fill.solid()
    lh.fill.fore_color.rgb = RGBColor(0xe7, 0x4c, 0x3c)
    lh.line.fill.background()
    lh_tf = lh.text_frame
    lh_tf.paragraphs[0].text = "⚠️ LIMITATION"
    lh_tf.paragraphs[0].font.size = Pt(16)
    lh_tf.paragraphs[0].font.bold = True
    lh_tf.paragraphs[0].font.color.rgb = WHITE
    lh_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Limitation content
    lc = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(5), Inches(4.5))
    tf = lc.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = limitation
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    
    # Arrow
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.2), Inches(3.8), Inches(0.9), Inches(0.5))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GREEN
    arrow.line.fill.background()
    
    # Solution box (right)
    sol_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(1.2), Inches(5.5), Inches(5.8))
    sol_box.fill.solid()
    sol_box.fill.fore_color.rgb = RGBColor(0xe8, 0xf8, 0xf5)  # Light green
    sol_box.line.color.rgb = GREEN
    sol_box.line.width = Pt(2)
    
    # Solution header
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.4), Inches(5.1), Inches(0.5))
    sh.fill.solid()
    sh.fill.fore_color.rgb = GREEN
    sh.line.fill.background()
    sh_tf = sh.text_frame
    sh_tf.paragraphs[0].text = f"✅ {solution_title}"
    sh_tf.paragraphs[0].font.size = Pt(16)
    sh_tf.paragraphs[0].font.bold = True
    sh_tf.paragraphs[0].font.color.rgb = WHITE
    sh_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Solution content
    sc = slide.shapes.add_textbox(Inches(7.6), Inches(2.1), Inches(5), Inches(4.5))
    tf = sc.text_frame
    tf.word_wrap = True
    for i, item in enumerate(solution_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
    
    return slide


# ============ CREATE THE PRESENTATION ============

# Slide 1: Title
add_title_slide(
    prs,
    "AI Chatbots vs LLMs vs Autonomous AI Agents",
    "Understanding the Evolution of Conversational AI Technologies"
)

# Slide 2: Agenda
add_content_slide(
    prs,
    "Agenda",
    [
        "What are AI Chatbots? – The Foundation",
        "What are Large Language Models (LLMs)? – The Brain",
        "What are Autonomous AI Agents? – The Complete Solution",
        "Evolution: How Each Technology Addresses Previous Limitations",
        "Real-World Applications & Use Cases",
        "Key Takeaways for Decision Makers"
    ],
    ACCENT_BLUE
)

# Slide 3: Evolution Overview
add_evolution_slide(
    prs,
    "The Evolution of Conversational AI",
    [
        ("AI CHATBOTS", "Rule-based systems with predefined responses. Limited to scripted conversations and keyword matching.", "🤖"),
        ("LLMs", "Deep learning models trained on vast text data. Generate human-like responses with contextual understanding.", "🧠"),
        ("AI AGENTS", "Autonomous systems that can reason, plan, use tools, and take actions to achieve complex goals.", "🚀")
    ]
)

# Slide 4: AI Chatbots - Definition
add_content_slide(
    prs,
    "AI Chatbots: The Foundation",
    [
        "Rule-based or simple ML systems designed for specific tasks",
        "Operate on predefined scripts, decision trees, and keyword matching",
        "Follow If-Then-Else logic patterns",
        "Popular examples: Traditional customer service bots, FAQ bots",
        "Technology: Pattern matching, regular expressions, basic NLP",
        "Best for: Simple, repetitive queries with known answers"
    ],
    RGBColor(0x95, 0xa5, 0xa6)
)

# Slide 5: Chatbot Limitations → LLM Solution
add_limitation_solution_slide(
    prs,
    "From Chatbots to LLMs: Addressing Key Limitations",
    "CHATBOT LIMITATIONS:\n\n• Cannot handle queries outside programmed scripts\n\n• Require manual updates for new scenarios\n\n• No understanding of context or nuance\n\n• Struggle with typos, slang, or variations\n\n• Cannot generate creative or novel responses\n\n• Limited to single-turn interactions",
    "LLMs SOLVE THIS",
    [
        "Understand natural language variations",
        "Generate contextually relevant responses",
        "Handle ambiguous or complex queries",
        "Learn patterns from vast training data",
        "Provide creative, nuanced answers",
        "Maintain conversation context"
    ],
    RGBColor(0xe7, 0x4c, 0x3c),
    GREEN
)

# Slide 6: LLMs - Definition
add_content_slide(
    prs,
    "Large Language Models (LLMs): The Brain",
    [
        "Deep neural networks trained on billions of text parameters",
        "Understand context, semantics, and generate human-like text",
        "Examples: GPT-4, Claude, Gemini, LLaMA, PaLM",
        "Can summarize, translate, code, reason, and create content",
        "Stateless by default – each conversation starts fresh",
        "Technology: Transformer architecture, attention mechanisms"
    ],
    ACCENT_BLUE
)

# Slide 7: LLM Capabilities
add_comparison_slide(
    prs,
    "LLM Capabilities vs Chatbots",
    "Traditional Chatbots",
    [
        "Keyword matching only",
        "Pre-written responses",
        "Single domain expertise",
        "Cannot handle variations",
        "No reasoning ability",
        "Static knowledge base"
    ],
    "Large Language Models",
    [
        "Semantic understanding",
        "Dynamic text generation",
        "Multi-domain knowledge",
        "Handle typos & variations",
        "Basic reasoning & logic",
        "Vast trained knowledge"
    ],
    RGBColor(0x95, 0xa5, 0xa6),
    ACCENT_BLUE
)

# Slide 8: LLM Limitations → Agent Solution
add_limitation_solution_slide(
    prs,
    "From LLMs to Agents: Bridging Critical Gaps",
    "LLM LIMITATIONS:\n\n• Cannot take real-world actions\n\n• No memory of past interactions\n\n• Cannot access real-time information\n\n• Single response only – no planning\n\n• Cannot use external tools or APIs\n\n• Knowledge cutoff date limitation",
    "AI AGENTS SOLVE THIS",
    [
        "Execute multi-step workflows",
        "Persistent memory & context",
        "Real-time web/database access",
        "Plan, reason, and iterate",
        "Integrate with any tool/API",
        "Always up-to-date information"
    ],
    RGBColor(0xe7, 0x4c, 0x3c),
    GREEN
)

# Slide 9: The Missing Piece - Tool/Function Calling
def add_tool_calling_slide(prs):
    """Special slide explaining tool calling as the bridge from LLM to Agent"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # Top accent
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🔧 The Secret Sauce: Tool Calling (Function Calling)"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Main explanation box
    main_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.333), Inches(1.4))
    main_box.fill.solid()
    main_box.fill.fore_color.rgb = RGBColor(0xf0, 0xf7, 0xff)
    main_box.line.color.rgb = ACCENT_BLUE
    main_box.line.width = Pt(2)
    
    main_tf = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(12), Inches(1.2))
    tf = main_tf.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Tool Calling allows LLMs to interact with external systems by generating structured function calls instead of just text. This is what transforms a chatbot into an agent."
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Diagram: LLM in center, tools around it
    # Center LLM box
    llm_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(3.2), Inches(3), Inches(1.2))
    llm_box.fill.solid()
    llm_box.fill.fore_color.rgb = ACCENT_BLUE
    llm_box.line.fill.background()
    llm_tf = llm_box.text_frame
    llm_tf.paragraphs[0].text = "🧠 LLM Core"
    llm_tf.paragraphs[0].font.size = Pt(20)
    llm_tf.paragraphs[0].font.bold = True
    llm_tf.paragraphs[0].font.color.rgb = WHITE
    llm_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Tool boxes around it
    tools = [
        ("🌐 Web Browser", 0.8, 2.8),
        ("💻 Code Interpreter", 0.8, 4.4),
        ("🔍 Search APIs", 9.5, 2.8),
        ("📊 Data Analysis", 9.5, 4.4),
    ]
    
    for tool_name, x, y in tools:
        t_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.8), Inches(0.9))
        t_box.fill.solid()
        t_box.fill.fore_color.rgb = GREEN
        t_box.line.fill.background()
        t_tf = t_box.text_frame
        t_tf.paragraphs[0].text = tool_name
        t_tf.paragraphs[0].font.size = Pt(16)
        t_tf.paragraphs[0].font.bold = True
        t_tf.paragraphs[0].font.color.rgb = WHITE
        t_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Arrows from LLM to tools
    # Left arrows
    arr1 = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(3.8), Inches(3.1), Inches(1.2), Inches(0.35))
    arr1.fill.solid()
    arr1.fill.fore_color.rgb = ORANGE
    arr1.line.fill.background()
    
    arr2 = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(3.8), Inches(4.6), Inches(1.2), Inches(0.35))
    arr2.fill.solid()
    arr2.fill.fore_color.rgb = ORANGE
    arr2.line.fill.background()
    
    # Right arrows
    arr3 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.4), Inches(3.1), Inches(1.0), Inches(0.35))
    arr3.fill.solid()
    arr3.fill.fore_color.rgb = ORANGE
    arr3.line.fill.background()
    
    arr4 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.4), Inches(4.6), Inches(1.0), Inches(0.35))
    arr4.fill.solid()
    arr4.fill.fore_color.rgb = ORANGE
    arr4.line.fill.background()
    
    # Link reference at bottom
    link_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.5))
    tf = link_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📚 Learn more: platform.openai.com/docs/guides/function-calling"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

add_tool_calling_slide(prs)

# Slide 10: Real Tool Examples (Generic)
def add_tool_examples_slide(prs):
    """Slide showing real-world tool examples"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Popular Tool Examples: LLM + Tools = Power"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Tool cards
    tools_data = [
        ("🌐", "Web Browsing", "ChatGPT + Bing", "Search real-time info, read articles, access current data beyond training cutoff"),
        ("💻", "Code Interpreter", "ChatGPT + Python", "Execute code, analyze data, create charts, process files in real-time"),
        ("🎨", "Image Generation", "ChatGPT + DALL-E", "Create images from text descriptions, edit existing images"),
        ("🔌", "Custom APIs", "Any LLM + Your API", "Connect to databases, CRMs, internal tools, trigger workflows"),
    ]
    
    y_positions = [1.1, 2.7, 4.3, 5.9]
    
    for i, (icon, name, example, desc) in enumerate(tools_data):
        # Icon circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y_positions[i]), Inches(0.9), Inches(0.9))
        circle.fill.solid()
        circle.fill.fore_color.rgb = GREEN if i % 2 == 0 else ACCENT_BLUE
        circle.line.fill.background()
        
        icon_tf = slide.shapes.add_textbox(Inches(0.5), Inches(y_positions[i] + 0.15), Inches(0.9), Inches(0.7))
        tf = icon_tf.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.CENTER
        
        # Tool name
        name_box = slide.shapes.add_textbox(Inches(1.6), Inches(y_positions[i]), Inches(2.5), Inches(0.5))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        
        # Example
        ex_box = slide.shapes.add_textbox(Inches(1.6), Inches(y_positions[i] + 0.45), Inches(2.5), Inches(0.4))
        tf = ex_box.text_frame
        p = tf.paragraphs[0]
        p.text = example
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = ORANGE
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(4.3), Inches(y_positions[i] + 0.1), Inches(8.5), Inches(0.8))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
    
    return slide

add_tool_examples_slide(prs)

# Slide 10.5: Our Real Implementation - Custom Tools
def add_our_tools_slide(prs):
    """Slide showing our actual implementation tools"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xf5, 0xf5, 0xf5)
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🛠️ Case Study: Our Custom Agent Tools"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(12.333), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Real tools from our Calendar Insights Agent (GPT-4.1-mini + OpenAI Function Calling)"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Tool 1: get_gdp
    tool1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.5), Inches(4), Inches(2.4))
    tool1.fill.solid()
    tool1.fill.fore_color.rgb = WHITE
    tool1.line.color.rgb = ACCENT_BLUE
    tool1.line.width = Pt(2)
    
    t1_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(3.8), Inches(0.45))
    t1_header.fill.solid()
    t1_header.fill.fore_color.rgb = ACCENT_BLUE
    t1_header.line.fill.background()
    t1h_tf = t1_header.text_frame
    t1h_tf.paragraphs[0].text = "🌍 get_gdp"
    t1h_tf.paragraphs[0].font.size = Pt(16)
    t1h_tf.paragraphs[0].font.bold = True
    t1h_tf.paragraphs[0].font.color.rgb = WHITE
    t1h_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    t1_content = slide.shapes.add_textbox(Inches(0.6), Inches(2.15), Inches(3.6), Inches(1.6))
    tf = t1_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "External API Tool"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p2 = tf.add_paragraph()
    p2.text = "• Calls api-ninjas.com/v1/gdp"
    p2.font.size = Pt(11)
    p2.font.color.rgb = DARK_GRAY
    p3 = tf.add_paragraph()
    p3.text = "• Params: country, year"
    p3.font.size = Pt(11)
    p3.font.color.rgb = DARK_GRAY
    p4 = tf.add_paragraph()
    p4.text = "• Returns real-time GDP data"
    p4.font.size = Pt(11)
    p4.font.color.rgb = DARK_GRAY
    
    # Tool 2: schedule_meeting
    tool2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.65), Inches(1.5), Inches(4), Inches(2.4))
    tool2.fill.solid()
    tool2.fill.fore_color.rgb = WHITE
    tool2.line.color.rgb = GREEN
    tool2.line.width = Pt(2)
    
    t2_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.75), Inches(1.6), Inches(3.8), Inches(0.45))
    t2_header.fill.solid()
    t2_header.fill.fore_color.rgb = GREEN
    t2_header.line.fill.background()
    t2h_tf = t2_header.text_frame
    t2h_tf.paragraphs[0].text = "📅 schedule_meeting"
    t2h_tf.paragraphs[0].font.size = Pt(16)
    t2h_tf.paragraphs[0].font.bold = True
    t2h_tf.paragraphs[0].font.color.rgb = WHITE
    t2h_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    t2_content = slide.shapes.add_textbox(Inches(4.85), Inches(2.15), Inches(3.6), Inches(1.6))
    tf = t2_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Action Tool (POST API)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p2 = tf.add_paragraph()
    p2.text = "• Creates calendar blocks"
    p2.font.size = Pt(11)
    p2.font.color.rgb = DARK_GRAY
    p3 = tf.add_paragraph()
    p3.text = "• ISO-8601 date/time params"
    p3.font.size = Pt(11)
    p3.font.color.rgb = DARK_GRAY
    p4 = tf.add_paragraph()
    p4.text = "• Connects to BriefingIQ API"
    p4.font.size = Pt(11)
    p4.font.color.rgb = DARK_GRAY
    
    # Tool 3: query_database
    tool3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(1.5), Inches(4), Inches(2.4))
    tool3.fill.solid()
    tool3.fill.fore_color.rgb = WHITE
    tool3.line.color.rgb = ORANGE
    tool3.line.width = Pt(2)
    
    t3_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(1.6), Inches(3.8), Inches(0.45))
    t3_header.fill.solid()
    t3_header.fill.fore_color.rgb = ORANGE
    t3_header.line.fill.background()
    t3h_tf = t3_header.text_frame
    t3h_tf.paragraphs[0].text = "🗄️ query_database"
    t3h_tf.paragraphs[0].font.size = Pt(16)
    t3h_tf.paragraphs[0].font.bold = True
    t3h_tf.paragraphs[0].font.color.rgb = WHITE
    t3h_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    t3_content = slide.shapes.add_textbox(Inches(9.1), Inches(2.15), Inches(3.6), Inches(1.6))
    tf = t3_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "NL-to-SQL Tool"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p2 = tf.add_paragraph()
    p2.text = "• Natural language → SQL"
    p2.font.size = Pt(11)
    p2.font.color.rgb = DARK_GRAY
    p3 = tf.add_paragraph()
    p3.text = "• Queries Oracle DB views"
    p3.font.size = Pt(11)
    p3.font.color.rgb = DARK_GRAY
    p4 = tf.add_paragraph()
    p4.text = "• Operations, Attendees, Revenue"
    p4.font.size = Pt(11)
    p4.font.color.rgb = DARK_GRAY
    
    # query_database detail box
    db_detail = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(4.1), Inches(12.5), Inches(3.0))
    db_detail.fill.solid()
    db_detail.fill.fore_color.rgb = WHITE
    db_detail.line.color.rgb = DARK_BLUE
    db_detail.line.width = Pt(2)
    
    db_header = slide.shapes.add_textbox(Inches(0.6), Inches(4.2), Inches(12), Inches(0.4))
    tf = db_header.text_frame
    p = tf.paragraphs[0]
    p.text = "🔍 query_database: Natural Language → SQL → Results"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Three Oracle views
    views = [
        ("VW_OPERATIONS_REPORT", "Event metadata, scheduling, logistics, ownership, regions", ACCENT_BLUE),
        ("VW_ATTENDEE_REPORT", "Attendee roster, decision makers, influencers, remote/in-person", GREEN),
        ("VW_OPP_TRACKING_REPORT", "Revenue, pipeline, opportunity metrics, probability of close", ORANGE),
    ]
    
    x_pos = [0.6, 4.4, 8.4]
    for i, (view_name, view_desc, color) in enumerate(views):
        v_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos[i]), Inches(4.75), Inches(3.8), Inches(1.1))
        v_box.fill.solid()
        v_box.fill.fore_color.rgb = color
        v_box.line.fill.background()
        
        v_tf = slide.shapes.add_textbox(Inches(x_pos[i] + 0.1), Inches(4.85), Inches(3.6), Inches(0.9))
        tf = v_tf.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = view_name
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p2 = tf.add_paragraph()
        p2.text = view_desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = WHITE
    
    # Example query
    ex_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.0), Inches(12), Inches(0.9))
    tf = ex_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = 'Example: "How many decision makers does Ford Motor have across all events?"'
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = DARK_GRAY
    p2 = tf.add_paragraph()
    p2.text = "→ LLM generates Oracle SQL → Executes on VW_ATTENDEE_REPORT → Returns structured data"
    p2.font.size = Pt(12)
    p2.font.color.rgb = ACCENT_BLUE
    
    return slide

add_our_tools_slide(prs)

# Slide 11: How Function Calling Works
def add_function_flow_slide(prs):
    """Show the flow of function calling"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "How Function Calling Works"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Flow steps
    steps = [
        ("1️⃣", "User Request", "\"What's the weather in NYC?\"", ACCENT_BLUE),
        ("2️⃣", "LLM Decides", "I need to call get_weather()", ACCENT_BLUE),
        ("3️⃣", "Tool Executes", "API returns: 72°F, Sunny", GREEN),
        ("4️⃣", "LLM Responds", "\"It's 72°F and sunny in NYC!\"", GREEN),
    ]
    
    x_positions = [0.4, 3.4, 6.5, 9.6]
    
    for i, (num, title, content, color) in enumerate(steps):
        # Step box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_positions[i]), Inches(1.3), Inches(2.9), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        
        # Number
        num_tf = slide.shapes.add_textbox(Inches(x_positions[i]), Inches(1.4), Inches(2.9), Inches(0.5))
        tf = num_tf.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER
        
        # Title
        t_tf = slide.shapes.add_textbox(Inches(x_positions[i]), Inches(1.9), Inches(2.9), Inches(0.5))
        tf = t_tf.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Content
        c_tf = slide.shapes.add_textbox(Inches(x_positions[i] + 0.1), Inches(2.45), Inches(2.7), Inches(0.9))
        tf = c_tf.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Arrow between steps
        if i < 3:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x_positions[i] + 2.95), Inches(2.2), Inches(0.4), Inches(0.3))
            arr.fill.solid()
            arr.fill.fore_color.rgb = ORANGE
            arr.line.fill.background()
    
    # Key insight box at bottom
    insight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.0), Inches(12.333), Inches(1.6))
    insight.fill.solid()
    insight.fill.fore_color.rgb = RGBColor(0xff, 0xf8, 0xe7)
    insight.line.color.rgb = ORANGE
    insight.line.width = Pt(2)
    
    insight_tf = slide.shapes.add_textbox(Inches(0.7), Inches(4.15), Inches(12), Inches(1.4))
    tf = insight_tf.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "💡 KEY INSIGHT"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    p2 = tf.add_paragraph()
    p2.text = "The LLM doesn't execute code or call APIs itself. It generates a structured request (JSON), the system executes it, and returns results to the LLM to formulate a response."
    p2.font.size = Pt(16)
    p2.font.color.rgb = DARK_GRAY
    
    # Before/After comparison at bottom
    before = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.9), Inches(5.8), Inches(1.3))
    before.fill.solid()
    before.fill.fore_color.rgb = RGBColor(0xff, 0xeb, 0xeb)
    before.line.color.rgb = RGBColor(0xe7, 0x4c, 0x3c)
    
    bef_tf = slide.shapes.add_textbox(Inches(0.7), Inches(6.0), Inches(5.4), Inches(1.1))
    tf = bef_tf.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "❌ Without Tools (Pure LLM)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xe7, 0x4c, 0x3c)
    p2 = tf.add_paragraph()
    p2.text = "\"I don't have access to real-time weather data.\""
    p2.font.size = Pt(13)
    p2.font.color.rgb = DARK_GRAY
    
    after = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(5.9), Inches(5.8), Inches(1.3))
    after.fill.solid()
    after.fill.fore_color.rgb = RGBColor(0xe8, 0xf8, 0xf5)
    after.line.color.rgb = GREEN
    
    aft_tf = slide.shapes.add_textbox(Inches(7.2), Inches(6.0), Inches(5.4), Inches(1.1))
    tf = aft_tf.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "✅ With Tools (Agent)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p2 = tf.add_paragraph()
    p2.text = "\"It's currently 72°F and sunny in New York City!\""
    p2.font.size = Pt(13)
    p2.font.color.rgb = DARK_GRAY
    
    return slide

add_function_flow_slide(prs)

# Slide 12: Technical Deep-Dive - Function Calling Code (Light theme, more readable)
def add_code_flow_slide(prs):
    """Technical slide showing the actual code/JSON for function calling"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Light background for readability
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xfa, 0xfa, 0xfc)
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.45))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🔧 Function Calling: The Technical Flow"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Code box colors - light theme
    code_bg = RGBColor(0x28, 0x2c, 0x34)  # VS Code dark
    
    # STEP 1: Define Tool (JSON Schema)
    step1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.15), Inches(0.6), Inches(4.25), Inches(2.85))
    step1.fill.solid()
    step1.fill.fore_color.rgb = code_bg
    step1.line.color.rgb = RGBColor(0x61, 0xaf, 0xef)  # Blue
    step1.line.width = Pt(3)
    
    s1_header = slide.shapes.add_textbox(Inches(0.25), Inches(0.65), Inches(4), Inches(0.35))
    tf = s1_header.text_frame
    p = tf.paragraphs[0]
    p.text = "1️⃣ Define Tool Schema"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x61, 0xaf, 0xef)
    
    s1_code = slide.shapes.add_textbox(Inches(0.25), Inches(1.0), Inches(4.1), Inches(2.35))
    tf = s1_code.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '''tools = [{
  "type": "function",
  "name": "get_gdp",
  "parameters": {
    "properties": {
      "country": {"type": "string"},
      "year": {"type": "string"}
    },
    "required": ["country","year"]
  }
}]'''
    p.font.size = Pt(10)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0x98, 0xc3, 0x79)  # Green
    
    # Arrow 1→2
    arr1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.43), Inches(1.9), Inches(0.4), Inches(0.3))
    arr1.fill.solid()
    arr1.fill.fore_color.rgb = RGBColor(0xe0, 0x6c, 0x75)  # Red accent
    arr1.line.fill.background()
    
    # STEP 2: API Call with Tools
    step2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.87), Inches(0.6), Inches(3.95), Inches(2.85))
    step2.fill.solid()
    step2.fill.fore_color.rgb = code_bg
    step2.line.color.rgb = RGBColor(0x98, 0xc3, 0x79)  # Green
    step2.line.width = Pt(3)
    
    s2_header = slide.shapes.add_textbox(Inches(4.97), Inches(0.65), Inches(3.8), Inches(0.35))
    tf = s2_header.text_frame
    p = tf.paragraphs[0]
    p.text = "2️⃣ Call API with Tools"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x98, 0xc3, 0x79)
    
    s2_code = slide.shapes.add_textbox(Inches(4.97), Inches(1.0), Inches(3.8), Inches(2.35))
    tf = s2_code.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '''response = client
  .responses.create(
    model="gpt-4.1-mini",
    tools=tools,
    input=[{
      "role": "user",
      "content": "USA GDP?"
    }]
)'''
    p.font.size = Pt(10)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0xe5, 0xc0, 0x7b)  # Yellow
    
    # Arrow 2→3
    arr2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.85), Inches(1.9), Inches(0.4), Inches(0.3))
    arr2.fill.solid()
    arr2.fill.fore_color.rgb = RGBColor(0xe0, 0x6c, 0x75)
    arr2.line.fill.background()
    
    # STEP 3: Model Returns Function Call
    step3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.28), Inches(0.6), Inches(3.85), Inches(2.85))
    step3.fill.solid()
    step3.fill.fore_color.rgb = code_bg
    step3.line.color.rgb = RGBColor(0xe0, 0x6c, 0x75)  # Red
    step3.line.width = Pt(3)
    
    s3_header = slide.shapes.add_textbox(Inches(9.38), Inches(0.65), Inches(3.6), Inches(0.35))
    tf = s3_header.text_frame
    p = tf.paragraphs[0]
    p.text = "3️⃣ Model Returns Call"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xe0, 0x6c, 0x75)
    
    s3_code = slide.shapes.add_textbox(Inches(9.38), Inches(1.0), Inches(3.6), Inches(2.35))
    tf = s3_code.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '''# response.output:
{
 "type":"function_call",
 "name": "get_gdp",
 "call_id": "abc123",
 "arguments": {
   "country": "USA",
   "year": "2024"
 }
}'''
    p.font.size = Pt(10)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0xc6, 0x78, 0xdd)  # Purple
    
    # STEP 4: Execute Tool (bottom left)
    step4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.15), Inches(3.6), Inches(4.25), Inches(1.7))
    step4.fill.solid()
    step4.fill.fore_color.rgb = code_bg
    step4.line.color.rgb = RGBColor(0xd1, 0x9a, 0x66)  # Orange
    step4.line.width = Pt(3)
    
    s4_header = slide.shapes.add_textbox(Inches(0.25), Inches(3.65), Inches(4), Inches(0.35))
    tf = s4_header.text_frame
    p = tf.paragraphs[0]
    p.text = "4️⃣ Execute Your Function"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xd1, 0x9a, 0x66)
    
    s4_code = slide.shapes.add_textbox(Inches(0.25), Inches(4.0), Inches(4.1), Inches(1.2))
    tf = s4_code.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '''args = json.loads(item.arguments)
result = get_gdp(args["country"],
                 args["year"])
# → {"gdp": 29167.78}'''
    p.font.size = Pt(10)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0x98, 0xc3, 0x79)
    
    # Arrow 4→5
    arr4 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.43), Inches(4.35), Inches(0.4), Inches(0.3))
    arr4.fill.solid()
    arr4.fill.fore_color.rgb = RGBColor(0xe0, 0x6c, 0x75)
    arr4.line.fill.background()
    
    # STEP 5: Send Tool Output Back
    step5 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.87), Inches(3.6), Inches(3.95), Inches(1.7))
    step5.fill.solid()
    step5.fill.fore_color.rgb = code_bg
    step5.line.color.rgb = RGBColor(0x56, 0xb6, 0xc2)  # Cyan
    step5.line.width = Pt(3)
    
    s5_header = slide.shapes.add_textbox(Inches(4.97), Inches(3.65), Inches(3.8), Inches(0.35))
    tf = s5_header.text_frame
    p = tf.paragraphs[0]
    p.text = "5️⃣ Append Tool Output"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x56, 0xb6, 0xc2)
    
    s5_code = slide.shapes.add_textbox(Inches(4.97), Inches(4.0), Inches(3.8), Inches(1.2))
    tf = s5_code.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '''input_list.append({
  "type":"function_call_output",
  "call_id": "abc123",
  "output": json.dumps(result)
})'''
    p.font.size = Pt(10)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0x56, 0xb6, 0xc2)
    
    # Arrow 5→6
    arr5 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.85), Inches(4.35), Inches(0.4), Inches(0.3))
    arr5.fill.solid()
    arr5.fill.fore_color.rgb = RGBColor(0xe0, 0x6c, 0x75)
    arr5.line.fill.background()
    
    # STEP 6: Second API Call → Final Response
    step6 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.28), Inches(3.6), Inches(3.85), Inches(1.7))
    step6.fill.solid()
    step6.fill.fore_color.rgb = code_bg
    step6.line.color.rgb = RGBColor(0xe5, 0xc0, 0x7b)  # Gold
    step6.line.width = Pt(3)
    
    s6_header = slide.shapes.add_textbox(Inches(9.38), Inches(3.65), Inches(3.6), Inches(0.35))
    tf = s6_header.text_frame
    p = tf.paragraphs[0]
    p.text = "6️⃣ Loop → Final Answer"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xe5, 0xc0, 0x7b)
    
    s6_code = slide.shapes.add_textbox(Inches(9.38), Inches(4.0), Inches(3.6), Inches(1.2))
    tf = s6_code.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '''# Call API again w/ output
response = client.responses
  .create(..., input=input_list)
# → "USA GDP: $29,167B"'''
    p.font.size = Pt(10)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0xe5, 0xc0, 0x7b)
    
    # Key insight at bottom - light theme
    insight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.15), Inches(5.5), Inches(12.98), Inches(1.85))
    insight.fill.solid()
    insight.fill.fore_color.rgb = RGBColor(0xff, 0xf3, 0xcd)  # Light yellow
    insight.line.color.rgb = RGBColor(0xff, 0x9f, 0x1c)
    insight.line.width = Pt(3)
    
    ins_tf = slide.shapes.add_textbox(Inches(0.35), Inches(5.6), Inches(12.6), Inches(1.7))
    tf = ins_tf.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "💡 THE AGENTIC LOOP"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xcc, 0x5a, 0x00)
    p2 = tf.add_paragraph()
    p2.text = "while True:  Send query + tools → If function_call: execute → append output → loop"
    p2.font.size = Pt(13)
    p2.font.name = "Consolas"
    p2.font.color.rgb = DARK_GRAY
    p3 = tf.add_paragraph()
    p3.text = "                       If text response: break → return final answer to user"
    p3.font.size = Pt(13)
    p3.font.name = "Consolas"
    p3.font.color.rgb = RGBColor(0x2e, 0x86, 0xab)
    
    return slide

add_code_flow_slide(prs)

# Slide 12.5: NL-to-SQL Flow (sqlite_qa.py)
def add_nl_to_sql_flow_slide(prs):
    """Technical slide showing the NL → SQL → Results flow"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Light background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xf8, 0xf9, 0xfa)
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.45))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🗄️ NL-to-SQL Tool: query_database Flow"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    code_bg = RGBColor(0x28, 0x2c, 0x34)
    
    # STEP 1: User Question
    step1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.15), Inches(0.6), Inches(3.15), Inches(2.0))
    step1.fill.solid()
    step1.fill.fore_color.rgb = RGBColor(0xe3, 0xf2, 0xfd)  # Light blue
    step1.line.color.rgb = ACCENT_BLUE
    step1.line.width = Pt(3)
    
    s1_header = slide.shapes.add_textbox(Inches(0.25), Inches(0.65), Inches(3), Inches(0.35))
    tf = s1_header.text_frame
    p = tf.paragraphs[0]
    p.text = "1️⃣ Natural Language"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    s1_content = slide.shapes.add_textbox(Inches(0.25), Inches(1.0), Inches(2.95), Inches(1.5))
    tf = s1_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '"How many decision makers does Ford have?"'
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = DARK_GRAY
    p2 = tf.add_paragraph()
    p2.text = "\n→ ask_sqlite(question)"
    p2.font.size = Pt(11)
    p2.font.name = "Consolas"
    p2.font.color.rgb = RGBColor(0x2e, 0x7d, 0x32)
    
    # Arrow 1→2
    arr1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.33), Inches(1.5), Inches(0.35), Inches(0.25))
    arr1.fill.solid()
    arr1.fill.fore_color.rgb = ORANGE
    arr1.line.fill.background()
    
    # STEP 2: Load Schema + Context
    step2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.72), Inches(0.6), Inches(3.15), Inches(2.0))
    step2.fill.solid()
    step2.fill.fore_color.rgb = code_bg
    step2.line.color.rgb = RGBColor(0xab, 0x47, 0xbc)  # Purple
    step2.line.width = Pt(3)
    
    s2_header = slide.shapes.add_textbox(Inches(3.82), Inches(0.65), Inches(3), Inches(0.35))
    tf = s2_header.text_frame
    p = tf.paragraphs[0]
    p.text = "2️⃣ Load Schema"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xce, 0x93, 0xd8)
    
    s2_code = slide.shapes.add_textbox(Inches(3.82), Inches(1.0), Inches(3), Inches(1.5))
    tf = s2_code.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '''# _load_schema()
VIEW_NAMES = [
  "VW_OPERATIONS_REPORT",
  "VW_ATTENDEE_REPORT",
  "VW_OPP_TRACKING_REPORT"
]'''
    p.font.size = Pt(9)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0xce, 0x93, 0xd8)
    
    # Arrow 2→3
    arr2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.9), Inches(1.5), Inches(0.35), Inches(0.25))
    arr2.fill.solid()
    arr2.fill.fore_color.rgb = ORANGE
    arr2.line.fill.background()
    
    # STEP 3: LLM Generates SQL
    step3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.28), Inches(0.6), Inches(3.15), Inches(2.0))
    step3.fill.solid()
    step3.fill.fore_color.rgb = code_bg
    step3.line.color.rgb = RGBColor(0x66, 0xbb, 0x6a)  # Green
    step3.line.width = Pt(3)
    
    s3_header = slide.shapes.add_textbox(Inches(7.38), Inches(0.65), Inches(3), Inches(0.35))
    tf = s3_header.text_frame
    p = tf.paragraphs[0]
    p.text = "3️⃣ LLM → SQL"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x66, 0xbb, 0x6a)
    
    s3_code = slide.shapes.add_textbox(Inches(7.38), Inches(1.0), Inches(3), Inches(1.5))
    tf = s3_code.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '''# _generate_sql()
response = client.responses
  .create(model=...,
    input=[schema+question])
→ {"sql": "SELECT..."}'''
    p.font.size = Pt(9)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0xa5, 0xd6, 0xa7)
    
    # Arrow 3→4
    arr3 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(10.46), Inches(1.5), Inches(0.35), Inches(0.25))
    arr3.fill.solid()
    arr3.fill.fore_color.rgb = ORANGE
    arr3.line.fill.background()
    
    # STEP 4: Execute SQL
    step4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.85), Inches(0.6), Inches(2.3), Inches(2.0))
    step4.fill.solid()
    step4.fill.fore_color.rgb = RGBColor(0xff, 0xf3, 0xe0)  # Light orange
    step4.line.color.rgb = ORANGE
    step4.line.width = Pt(3)
    
    s4_header = slide.shapes.add_textbox(Inches(10.95), Inches(0.65), Inches(2.1), Inches(0.35))
    tf = s4_header.text_frame
    p = tf.paragraphs[0]
    p.text = "4️⃣ Execute"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    
    s4_code = slide.shapes.add_textbox(Inches(10.95), Inches(1.0), Inches(2.1), Inches(1.5))
    tf = s4_code.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '''# Oracle DB
conn.execute(
  text(sql)
)
→ rows'''
    p.font.size = Pt(9)
    p.font.name = "Consolas"
    p.font.color.rgb = DARK_GRAY
    
    # Generated SQL Box (spanning bottom)
    sql_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.15), Inches(2.8), Inches(12.98), Inches(1.5))
    sql_box.fill.solid()
    sql_box.fill.fore_color.rgb = code_bg
    sql_box.line.color.rgb = RGBColor(0x42, 0xa5, 0xf5)
    sql_box.line.width = Pt(2)
    
    sql_header = slide.shapes.add_textbox(Inches(0.3), Inches(2.85), Inches(12.5), Inches(0.35))
    tf = sql_header.text_frame
    p = tf.paragraphs[0]
    p.text = "📝 Generated SQL (Oracle syntax with rules applied):"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x42, 0xa5, 0xf5)
    
    sql_code = slide.shapes.add_textbox(Inches(0.3), Inches(3.2), Inches(12.5), Inches(1.0))
    tf = sql_code.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '''SELECT CUSTOMERNAME, COUNT(*) as decision_maker_count FROM VW_ATTENDEE_REPORT
WHERE lower(CUSTOMERNAME) LIKE '%ford%' AND DECISIONMAKER = 'Yes' GROUP BY CUSTOMERNAME FETCH FIRST 100 ROWS ONLY'''
    p.font.size = Pt(11)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0xe5, 0xc0, 0x7b)
    
    # Oracle Rules Box
    rules_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.15), Inches(4.45), Inches(6.4), Inches(2.85))
    rules_box.fill.solid()
    rules_box.fill.fore_color.rgb = RGBColor(0xff, 0xeb, 0xee)  # Light red
    rules_box.line.color.rgb = RGBColor(0xef, 0x53, 0x50)
    rules_box.line.width = Pt(2)
    
    rules_header = slide.shapes.add_textbox(Inches(0.3), Inches(4.5), Inches(6.2), Inches(0.35))
    tf = rules_header.text_frame
    p = tf.paragraphs[0]
    p.text = "⚠️ Oracle SQL Rules (in system prompt):"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xc6, 0x28, 0x28)
    
    rules_content = slide.shapes.add_textbox(Inches(0.3), Inches(4.85), Inches(6.2), Inches(2.3))
    tf = rules_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "• NEVER end with semicolon"
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_GRAY
    p2 = tf.add_paragraph()
    p2.text = "• Use FETCH FIRST n ROWS ONLY (not LIMIT)"
    p2.font.size = Pt(11)
    p2.font.color.rgb = DARK_GRAY
    p3 = tf.add_paragraph()
    p3.text = "• Epoch ms → date: date '1970-01-01' + (ms/1000)/86400"
    p3.font.size = Pt(11)
    p3.font.color.rgb = DARK_GRAY
    p4 = tf.add_paragraph()
    p4.text = "• Case-insensitive: lower(col) LIKE '%term%'"
    p4.font.size = Pt(11)
    p4.font.color.rgb = DARK_GRAY
    p5 = tf.add_paragraph()
    p5.text = "• Only SELECT queries allowed"
    p5.font.size = Pt(11)
    p5.font.color.rgb = DARK_GRAY
    
    # Return Value Box
    return_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.73), Inches(4.45), Inches(6.4), Inches(2.85))
    return_box.fill.solid()
    return_box.fill.fore_color.rgb = RGBColor(0xe8, 0xf5, 0xe9)  # Light green
    return_box.line.color.rgb = GREEN
    return_box.line.width = Pt(2)
    
    return_header = slide.shapes.add_textbox(Inches(6.88), Inches(4.5), Inches(6.2), Inches(0.35))
    tf = return_header.text_frame
    p = tf.paragraphs[0]
    p.text = "✅ Return Value (back to agent):"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x2e, 0x7d, 0x32)
    
    return_code = slide.shapes.add_textbox(Inches(6.88), Inches(4.85), Inches(6.2), Inches(2.3))
    tf = return_code.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '''{
  "sql": "SELECT CUSTOMERNAME...",
  "explanation": "Query counts...",
  "columns": ["CUSTOMERNAME", 
              "DECISION_MAKER_COUNT"],
  "rows": [
    {"CUSTOMERNAME": "Ford Motor",
     "DECISION_MAKER_COUNT": 12}
  ]
}'''
    p.font.size = Pt(10)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0x2e, 0x7d, 0x32)
    
    return slide

add_nl_to_sql_flow_slide(prs)

# Slide 13: From Tools to Agents
add_content_slide(
    prs,
    "From Tool-Enabled LLMs to Autonomous Agents",
    [
        "LLM + 1 Tool = Enhanced Capability (e.g., ChatGPT + Web)",
        "LLM + Multiple Tools = Flexible Problem Solving",
        "LLM + Tools + Memory = Contextual Assistant",
        "LLM + Tools + Memory + Planning = AUTONOMOUS AGENT",
        "Agent = LLM that can decide WHICH tools to use, WHEN to use them, and iterate until goal is achieved"
    ],
    GREEN
)

# Slide 13: Three-way comparison (condensed)
add_three_column_slide(
    prs,
    "Complete Comparison: Chatbots vs LLMs vs AI Agents",
    ("AI CHATBOTS", [
        "Rule-based logic",
        "Scripted responses",
        "No external access",
        "Cannot take actions",
        "No memory",
        "Simple queries only"
    ]),
    ("LLMs", [
        "Pattern recognition",
        "Generated responses",
        "Knowledge from training",
        "Text output only",
        "Session-based context",
        "Complex conversations"
    ]),
    ("AI AGENTS", [
        "Reasoning + planning",
        "Tool-enabled actions",
        "Real-time data access",
        "Execute workflows",
        "Persistent memory",
        "Goal-oriented behavior"
    ])
)

# Slide 14: When to Use What
add_content_slide(
    prs,
    "Choosing the Right Technology",
    [
        "CHATBOTS: Simple, repetitive, high-volume tasks with known answers",
        "LLMs: Language understanding, content generation, flexible conversations",
        "LLM + Tools: When you need real-time data or specific capabilities",
        "AI AGENTS: Complex goals requiring planning, multiple tools, and iteration",
        "Start simple → Add tools → Scale to agents as needed"
    ],
    ORANGE
)

# Slide 15: Key Takeaways
add_content_slide(
    prs,
    "Key Takeaways",
    [
        "CHATBOTS → LLMs: From scripts to understanding",
        "LLMs → LLM+Tools: From knowing to accessing (via Function Calling)",
        "LLM+Tools → Agents: From single actions to autonomous goal pursuit",
        "TOOL CALLING is the bridge that enables LLMs to interact with the world",
        "Agents = LLM + Tools + Memory + Planning + Autonomy"
    ],
    DARK_BLUE
)

# Slide 15: Closing
add_title_slide(
    prs,
    "Questions?",
    "From Chatbots to Agents: The AI Evolution Continues"
)

# Save the presentation
output_path = "/Users/jyothiraditya/projects-main/calender-insights/AI_Chatbots_LLMs_Agents_Comparison.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")

