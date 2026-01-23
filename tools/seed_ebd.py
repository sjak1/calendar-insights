"""
EBD Seeding Script - Fill template with dummy data for testing

Creates a filled-out EBD document from the template for AI testing purposes.
Run: python tools/seed_ebd.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path


# Dummy data for a fictional company
DUMMY_DATA = {
    "company_name": "Amazon",
    "industry": "Retail",
    "address": "410 Terry Avenue North, Seattle, WA 98109",
    "website": "www.amazon.com",
    
    "meeting_date": "February 10, 2026 | 10:00 AM - 5:00 PM PST",
    "meeting_location": "Oracle Executive Briefing Center - Redwood City",
    "engagement_type": "1:1 Executive Briefing",
    
    "ebd_contact": "Sarah Mitchell, Sr. Account Executive, sarah.mitchell@oracle.com, (512) 555-1234",
    
    "oracle_attendees": """Thomas Reed, VP Analytics Solutions, thomas.reed@oracle.com
Sarah Kim, Sr. Director Customer Experience, sarah.kim@oracle.com
James Liu, Principal Data Architect, james.liu@oracle.com""",
    
    "customer_attendees": """Jennifer Martinez, VP Retail Analytics
David Kim, Director Supply Chain Technology
Patricia Lee, Head of Customer Insights
Michael Brown, Senior Data Scientist""",
    
    "company_description": """NovaTech Industries is a $2.8B manufacturing company specializing in precision aerospace components and industrial automation systems. Founded in 1987, they operate 12 manufacturing facilities across North America and employ 8,500 people. They are a key supplier to Boeing, Lockheed Martin, and General Electric. NovaTech is currently undergoing a digital transformation initiative to modernize their legacy ERP systems and improve supply chain visibility.""",
    
    "meeting_objectives": """1. Understand how Oracle Cloud ERP can consolidate their 4 legacy ERP systems (SAP, JD Edwards, custom solutions) into a unified platform
2. Evaluate Oracle's manufacturing execution and supply chain planning capabilities for aerospace compliance requirements
3. Discuss AI/ML capabilities for predictive maintenance and demand forecasting
4. Review TCO comparison and implementation timeline for a phased rollout approach
5. Meet with Oracle executives to establish strategic partnership for their 5-year digital transformation roadmap""",
    
    "business_challenges": """1. LEGACY SYSTEM FRAGMENTATION: Running 4 different ERP systems across facilities, causing data silos, reconciliation delays, and $3M+ annually in integration maintenance costs

2. SUPPLY CHAIN VISIBILITY: Limited real-time visibility into supplier inventory and production capacity; recent chip shortage caused 6-week production delays and $12M in lost revenue

3. COMPLIANCE & TRACEABILITY: Aerospace customers (Boeing, Lockheed) requiring AS9100 compliance with full component traceability; current systems lack unified audit trail

4. WORKFORCE CHALLENGES: 40% of senior manufacturing engineers retiring in next 5 years; need to capture institutional knowledge and automate manual processes""",
    
    "it_strategy": """CEO Vision: "Become the most digitally advanced precision manufacturer in North America by 2028"

IT Strategy:
- Cloud-first approach for all new applications
- Consolidate to single ERP platform within 3 years
- Implement IoT sensors across all production lines for real-time monitoring
- Leverage AI for predictive maintenance (targeting 30% reduction in unplanned downtime)
- Modernize B2B integration with suppliers and customers""",
    
    "customer_lifecycle": """SALES STAGE: Evaluation / Active Opportunity ($8.5M TCV)
- Initial discovery completed Q3 2025
- RFP issued to Oracle, SAP, and Microsoft
- Technical deep-dive scheduled for February 2026
- Decision expected Q2 2026

WHY CONSIDERING ORACLE:
- Existing Oracle Database footprint (60% of data warehouses)
- Positive experience with Oracle support
- Industry references from aerospace peers (Honeywell, Textron)

RECENT IMPLEMENTATIONS:
- Oracle Analytics Cloud deployed in 2024 (successful)
- Oracle Autonomous Database POC completed (positive results)""",
    
    "account_status": """GREEN - Strong relationship

POSITIVES:
- Executive sponsor (CIO James Wilson) is former Oracle customer at previous company
- CFO impressed with Oracle's financial management capabilities in demos
- IT team prefers Oracle's integration capabilities

POTENTIAL DERAILERS:
- CEO has personal relationship with SAP executive
- Concerns about implementation timeline (want go-live within 18 months)
- Board pressure to show quick ROI""",
    
    "oracle_talking_points": """1. INDUSTRY EXPERTISE: Highlight Oracle's manufacturing cloud customers in aerospace (Honeywell, Textron, Spirit AeroSystems) and our AS9100 compliance capabilities

2. UNIFIED PLATFORM VALUE: Demonstrate how Oracle Fusion Cloud eliminates integration costs and provides single source of truth - reference $4M+ annual savings at similar-sized manufacturers

3. AI/INNOVATION LEADERSHIP: Showcase Oracle's embedded AI capabilities for demand sensing, predictive maintenance, and supply chain optimization - differentiate from SAP's bolt-on approach""",
    
    # Attendee profiles
    "attendee_1_name": "Robert Hayes, CEO, robert.hayes@novatech.com",
    "attendee_1_perspective": "Focused on shareholder value and competitive positioning. Wants to ensure NovaTech remains supplier of choice for aerospace OEMs. Concerned about transformation risk and business disruption. Has worked with SAP at previous company.",
    "attendee_1_bio": "30+ years in manufacturing, former COO at Precision Castparts. Stanford MBA. Board member at National Association of Manufacturers.",
    
    "attendee_2_name": "Amanda Foster, CFO, amanda.foster@novatech.com",
    "attendee_2_perspective": "Driving cost reduction and operational efficiency. Very analytical - will want detailed TCO and ROI projections. Interested in financial close automation and real-time reporting.",
    "attendee_2_bio": "Former CFO at Collins Aerospace. CPA with 20 years finance experience. Published author on manufacturing finance transformation.",
    
    "attendee_3_name": "James Wilson, CIO, james.wilson@novatech.com",
    "attendee_3_perspective": "Oracle advocate - used Oracle successfully at previous company. Main champion for this initiative. Concerned about implementation partner quality and timeline.",
    "attendee_3_bio": "Former VP of IT at Rockwell Collins. 25 years in manufacturing IT. Led successful Oracle implementation at previous employer.",
    
    # Customer references
    "ref_1_name": "Honeywell Aerospace",
    "ref_1_industry": "Aerospace & Defense",
    "ref_1_website": "www.honeywell.com",
    "ref_1_products": "Oracle Cloud ERP, Oracle SCM Cloud, Oracle Manufacturing Cloud",
    "ref_1_summary": "Consolidated 5 legacy ERP systems to Oracle Cloud, achieving 40% reduction in close time and $15M annual savings in IT costs. Full AS9100 compliance maintained throughout transformation.",
    
    "ref_2_name": "Spirit AeroSystems",
    "ref_2_industry": "Aerospace Manufacturing",
    "ref_2_website": "www.spiritaero.com",
    "ref_2_products": "Oracle Cloud ERP, Oracle EPM Cloud",
    "ref_2_summary": "Implemented Oracle for financial planning and manufacturing operations. Achieved real-time visibility across 10 facilities and reduced inventory carrying costs by 25%.",
    
    # Recent engagements
    "recent_engagement_1": "Technical Deep Dive | October 2025 / Virtual | Sarah Mitchell | James Wilson, Lisa Chang",
    "recent_engagement_2": "Oracle OpenWorld Visit | September 2025 / Las Vegas | Regional VP | Robert Hayes, Amanda Foster",
    
    # Account team
    "scd_lam": "Sarah Mitchell, Sr. Account Executive, sarah.mitchell@oracle.com, (512) 555-1234",
    "other_am": "Kevin O'Brien, Strategic Account Manager, kevin.obrien@oracle.com, (512) 555-5678",
    "cse_csm": "Rachel Torres, Customer Success Manager, rachel.torres@oracle.com, (512) 555-9012",
    "exec_sponsor": "Thomas Anderson, SVP North America Sales, thomas.anderson@oracle.com",
    "impl_partner": "Deloitte - Manufacturing Practice (Primary), Infosys (Secondary for integrations)",
    
    # LMS/SIA info
    "last_audit_date": "March 2024",
    "audit_active": "No",
    "lms_legal": "No",
    "ula_cert": "Yes - Expires Dec 2026",
    "lms_contact": "John Smith, LMS Specialist",
    "sia_engaged": "Yes",
    "sia_contact": "Maria Garcia, SIA Manager",
    
    # Third customer reference
    "ref_3_name": "Textron Aviation",
    "ref_3_industry": "Aerospace Manufacturing",
    "ref_3_website": "www.txtav.com",
    "ref_3_products": "Oracle Cloud ERP, Oracle IoT Cloud",
    "ref_3_summary": "Implemented Oracle Cloud for manufacturing operations across 5 plants. Achieved 35% improvement in production planning accuracy and real-time visibility into shop floor operations.",
    
    # Financial data (Slide 5)
    "total_deals_size": "$8.5M",
    "customer_segment": "Enterprise",
    "tam": "$45M",
    "avg_oracle_spend": "$2.1M",
    "annual_revenue": "$2.8B",
    "share_of_wallet": "4.7%",
    
    # Product footprint
    "apps_support_spend": "$0.8M",
    "apps_footprint": "JD Edwards (legacy)",
    "apps_competitor": "SAP S/4HANA",
    "db_support_spend": "$1.2M",
    "db_footprint": "Oracle Database 19c",
    "db_competitor": "Microsoft SQL Server",
    "middleware_support_spend": "$0.1M",
    "middleware_footprint": "WebLogic",
    "middleware_competitor": "IBM MQ",
    
    # Open opportunities
    "opp_cloud_saas": "$4.2M",
    "opp_cloud_saas_prob": "65%",
    "opp_cloud_paas": "$2.8M", 
    "opp_cloud_paas_prob": "55%",
    "opp_services": "$1.5M",
    "opp_services_prob": "70%",
    "recent_win_analytics": "$0.3M",
    "recent_win_adb": "$0.2M",
}


def fill_table_cell(table, row_idx, col_idx, text):
    """Safely fill a table cell if it exists."""
    try:
        if row_idx < len(table.rows) and col_idx < len(table.rows[row_idx].cells):
            cell = table.rows[row_idx].cells[col_idx]
            # Clear existing content
            cell.text = text
            return True
    except Exception as e:
        print(f"Warning: Could not fill cell [{row_idx}][{col_idx}]: {e}")
    return False


def find_and_replace_text(shape, old_text, new_text):
    """Find and replace text in a shape's text frame."""
    if not shape.has_text_frame:
        return False
    
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if old_text in run.text:
                run.text = run.text.replace(old_text, new_text)
                return True
    return False


def seed_ebd_document(template_path: str, output_path: str):
    """
    Fill the EBD template with dummy data.
    
    Args:
        template_path: Path to the template PPTX
        output_path: Path for the filled output PPTX
    """
    print(f"Loading template: {template_path}")
    prs = Presentation(template_path)
    
    print(f"Found {len(prs.slides)} slides")
    
    # Process each slide
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        print(f"\nProcessing Slide {slide_num}...")
        
        # Skip slide 1 (instructions)
        if slide_num == 1:
            print("  Skipping instructions slide")
            continue
        
        # Find all tables and text shapes
        tables = []
        text_shapes = []
        
        for shape in slide.shapes:
            if shape.has_table:
                tables.append(shape.table)
            if shape.has_text_frame:
                text_shapes.append(shape)
        
        print(f"  Found {len(tables)} tables, {len(text_shapes)} text shapes")
        
        # Slide 2: Main EBD form
        if slide_num == 2 and tables:
            main_table = tables[0]
            print(f"  Main table: {len(main_table.rows)} rows x {len(main_table.columns)} cols")
            
            # Fill in the main table based on position
            # Row 0: Company Name, Industry
            fill_table_cell(main_table, 0, 0, DUMMY_DATA["company_name"])
            fill_table_cell(main_table, 0, 2, DUMMY_DATA["industry"])
            
            # Row 1-6: Address, Meeting info (left side) and details (right side)
            fill_table_cell(main_table, 1, 0, DUMMY_DATA["address"])
            fill_table_cell(main_table, 1, 2, "Meeting Date/Time")
            fill_table_cell(main_table, 1, 3, DUMMY_DATA["meeting_date"])
            
            fill_table_cell(main_table, 2, 2, "Meeting Location")
            fill_table_cell(main_table, 2, 3, DUMMY_DATA["meeting_location"])
            
            fill_table_cell(main_table, 3, 0, DUMMY_DATA["address"])
            fill_table_cell(main_table, 3, 2, "Engagement Type")
            fill_table_cell(main_table, 3, 3, DUMMY_DATA["engagement_type"])
            
            fill_table_cell(main_table, 4, 0, DUMMY_DATA["website"])
            fill_table_cell(main_table, 4, 2, "EBD/Meeting Point of Contact")
            fill_table_cell(main_table, 4, 3, DUMMY_DATA["ebd_contact"])
            
            fill_table_cell(main_table, 5, 2, "Oracle Meeting Attendees")
            fill_table_cell(main_table, 5, 3, DUMMY_DATA["oracle_attendees"])
            
            fill_table_cell(main_table, 6, 2, "Customer Meeting Attendees")
            fill_table_cell(main_table, 6, 3, DUMMY_DATA["customer_attendees"])
            
            # Row 7: Company Description
            fill_table_cell(main_table, 7, 1, DUMMY_DATA["company_description"])
            
            # Rows 8-13: The key EBD fields
            fill_table_cell(main_table, 8, 1, DUMMY_DATA["meeting_objectives"])
            fill_table_cell(main_table, 9, 1, DUMMY_DATA["business_challenges"])
            fill_table_cell(main_table, 10, 1, DUMMY_DATA["it_strategy"])
            fill_table_cell(main_table, 11, 1, DUMMY_DATA["customer_lifecycle"])
            fill_table_cell(main_table, 12, 1, DUMMY_DATA["account_status"])
            fill_table_cell(main_table, 13, 1, DUMMY_DATA["oracle_talking_points"])
            
            print("  ✓ Filled main EBD table")
        
        # Slide 3: Account team and attendee profiles
        if slide_num == 3 and tables:
            print(f"  Found {len(tables)} tables on slide 3")
            
            for table_idx, table in enumerate(tables):
                rows = len(table.rows)
                cols = len(table.columns)
                print(f"    Table {table_idx + 1}: {rows}x{cols}")
                
                # Table 1 (11x5): Recent Engagements + Account Team
                if rows == 11 and cols == 5:
                    # Row 0: Header "Most Recent Oracle Executive Engagements"
                    # Row 1: Column headers
                    # Rows 2-3: Engagement data
                    fill_table_cell(table, 2, 0, "1:1 Executive Briefing")
                    fill_table_cell(table, 2, 1, "October 15, 2025 / Virtual")
                    fill_table_cell(table, 2, 2, "Sarah Mitchell")
                    fill_table_cell(table, 2, 3, "James Wilson, Lisa Chang")
                    
                    fill_table_cell(table, 3, 0, "Oracle OpenWorld")
                    fill_table_cell(table, 3, 1, "September 2025 / Las Vegas")
                    fill_table_cell(table, 3, 2, "Thomas Anderson, SVP")
                    fill_table_cell(table, 3, 3, "Robert Hayes, Amanda Foster")
                    
                    # Account Team section (rows 4-10)
                    # Find and fill based on label text
                    for row_idx in range(4, rows):
                        cell_text = table.rows[row_idx].cells[0].text.strip().lower()
                        if "scd" in cell_text or "lam" in cell_text:
                            fill_table_cell(table, row_idx, 1, DUMMY_DATA["scd_lam"])
                        elif "other account" in cell_text:
                            fill_table_cell(table, row_idx, 1, DUMMY_DATA["other_am"])
                        elif "cse" in cell_text or "csm" in cell_text:
                            fill_table_cell(table, row_idx, 1, DUMMY_DATA["cse_csm"])
                        elif "executive sponsor" in cell_text:
                            fill_table_cell(table, row_idx, 1, DUMMY_DATA["exec_sponsor"])
                        elif "implementation" in cell_text or "partner" in cell_text:
                            fill_table_cell(table, row_idx, 1, DUMMY_DATA["impl_partner"])
                    print("    ✓ Filled engagements and account team")
                
                # Table 2 (4x5): LMS/SIA info
                if rows == 4 and cols == 5:
                    # Row 0: Headers (Last Audit Date, Audit Active, LMS Legal, ULA Cert, LMS Contact)
                    # Row 1: Values
                    fill_table_cell(table, 1, 0, DUMMY_DATA["last_audit_date"])
                    fill_table_cell(table, 1, 1, DUMMY_DATA["audit_active"])
                    fill_table_cell(table, 1, 2, DUMMY_DATA["lms_legal"])
                    fill_table_cell(table, 1, 3, DUMMY_DATA["ula_cert"])
                    fill_table_cell(table, 1, 4, DUMMY_DATA["lms_contact"])
                    # Row 2: SIA Engaged header
                    # Row 3: SIA values
                    fill_table_cell(table, 3, 0, DUMMY_DATA["sia_engaged"])
                    fill_table_cell(table, 3, 4, DUMMY_DATA["sia_contact"])
                    print("    ✓ Filled LMS/SIA info")
                
                # Table 3 (9x3): Attendee profiles
                if rows == 9 and cols == 3:
                    # Each attendee takes 3 rows: Name, Perspective, Bio
                    # Column 1 has the labels, Column 2 has the values (or maybe cols 1-2 merged)
                    # Try filling in column 2 first, then column 1 if that doesn't work
                    
                    # Attendee 1 (rows 0-2)
                    fill_table_cell(table, 0, 1, DUMMY_DATA["attendee_1_name"])
                    fill_table_cell(table, 0, 2, DUMMY_DATA["attendee_1_name"])
                    fill_table_cell(table, 1, 1, DUMMY_DATA["attendee_1_perspective"])
                    fill_table_cell(table, 1, 2, DUMMY_DATA["attendee_1_perspective"])
                    fill_table_cell(table, 2, 1, DUMMY_DATA["attendee_1_bio"])
                    fill_table_cell(table, 2, 2, DUMMY_DATA["attendee_1_bio"])
                    
                    # Attendee 2 (rows 3-5)
                    fill_table_cell(table, 3, 1, DUMMY_DATA["attendee_2_name"])
                    fill_table_cell(table, 3, 2, DUMMY_DATA["attendee_2_name"])
                    fill_table_cell(table, 4, 1, DUMMY_DATA["attendee_2_perspective"])
                    fill_table_cell(table, 4, 2, DUMMY_DATA["attendee_2_perspective"])
                    fill_table_cell(table, 5, 1, DUMMY_DATA["attendee_2_bio"])
                    fill_table_cell(table, 5, 2, DUMMY_DATA["attendee_2_bio"])
                    
                    # Attendee 3 (rows 6-8)
                    fill_table_cell(table, 6, 1, DUMMY_DATA["attendee_3_name"])
                    fill_table_cell(table, 6, 2, DUMMY_DATA["attendee_3_name"])
                    fill_table_cell(table, 7, 1, DUMMY_DATA["attendee_3_perspective"])
                    fill_table_cell(table, 7, 2, DUMMY_DATA["attendee_3_perspective"])
                    fill_table_cell(table, 8, 1, DUMMY_DATA["attendee_3_bio"])
                    fill_table_cell(table, 8, 2, DUMMY_DATA["attendee_3_bio"])
                    print("    ✓ Filled attendee profiles")
        
        # Slide 4: Customer references
        if slide_num == 4 and tables:
            print(f"  Found {len(tables)} reference tables")
            
            for table_idx, table in enumerate(tables):
                rows = len(table.rows)
                cols = len(table.columns)
                print(f"    Table {table_idx + 1}: {rows}x{cols}")
                
                if rows == 6:  # Customer reference table structure (6x2)
                    ref_data = None
                    if table_idx == 0:
                        ref_data = ("ref_1_name", "ref_1_industry", "ref_1_website", "ref_1_products", "", "ref_1_summary")
                    elif table_idx == 1:
                        ref_data = ("ref_2_name", "ref_2_industry", "ref_2_website", "ref_2_products", "", "ref_2_summary")
                    elif table_idx == 2:
                        ref_data = ("ref_3_name", "ref_3_industry", "ref_3_website", "ref_3_products", "", "ref_3_summary")
                    
                    if ref_data:
                        for row_idx, key in enumerate(ref_data):
                            if key and key in DUMMY_DATA:
                                fill_table_cell(table, row_idx, 1, DUMMY_DATA[key])
                        print(f"    ✓ Filled reference {table_idx + 1}")
        
        # Slide 5: Financial data and product footprint
        if slide_num == 5 and tables:
            print(f"  Found {len(tables)} tables on slide 5")
            
            for table_idx, table in enumerate(tables):
                rows = len(table.rows)
                cols = len(table.columns)
                print(f"    Table {table_idx + 1}: {rows}x{cols}")
                
                # Main financial table (30x10)
                # Structure:
                # - Col 0: Product Pillar/Label
                # - Col 2: Product Line
                # - Col 4: Deal Size / Value
                # - Col 5: Probability
                # - Col 7: Recent Wins / Second label
                # - Col 9: Second value
                if rows >= 25 and cols >= 10:
                    # Row 8: SaaS opportunities
                    fill_table_cell(table, 8, 4, DUMMY_DATA["opp_cloud_saas"])
                    fill_table_cell(table, 8, 5, DUMMY_DATA["opp_cloud_saas_prob"])
                    
                    # Row 9: PaaS/IaaS opportunities
                    fill_table_cell(table, 9, 4, DUMMY_DATA["opp_cloud_paas"])
                    fill_table_cell(table, 9, 5, DUMMY_DATA["opp_cloud_paas_prob"])
                    
                    # Row 13: Services opportunities
                    fill_table_cell(table, 13, 4, DUMMY_DATA["opp_services"])
                    fill_table_cell(table, 13, 5, DUMMY_DATA["opp_services_prob"])
                    
                    # Row 14: Total Deals Size | Customer Segment
                    fill_table_cell(table, 14, 4, DUMMY_DATA["total_deals_size"])
                    fill_table_cell(table, 14, 9, DUMMY_DATA["customer_segment"])
                    
                    # Row 15: TAM | Avg Annual Oracle Spend
                    fill_table_cell(table, 15, 4, DUMMY_DATA["tam"])
                    fill_table_cell(table, 15, 9, DUMMY_DATA["avg_oracle_spend"])
                    
                    # Row 16: Annual Company Revenue | Share of Wallet
                    fill_table_cell(table, 16, 4, DUMMY_DATA["annual_revenue"])
                    fill_table_cell(table, 16, 9, DUMMY_DATA["share_of_wallet"])
                    
                    # Product Footprint section (rows 19-29)
                    # Row 19: Applications (under Software & Cloud)
                    fill_table_cell(table, 19, 3, DUMMY_DATA["apps_support_spend"])
                    fill_table_cell(table, 19, 5, DUMMY_DATA["apps_footprint"])
                    fill_table_cell(table, 19, 7, DUMMY_DATA["apps_competitor"])
                    
                    # Row 20: Database
                    fill_table_cell(table, 20, 3, DUMMY_DATA["db_support_spend"])
                    fill_table_cell(table, 20, 5, DUMMY_DATA["db_footprint"])
                    fill_table_cell(table, 20, 7, DUMMY_DATA["db_competitor"])
                    
                    # Row 21: Middleware
                    fill_table_cell(table, 21, 3, DUMMY_DATA["middleware_support_spend"])
                    fill_table_cell(table, 21, 5, DUMMY_DATA["middleware_footprint"])
                    fill_table_cell(table, 21, 7, DUMMY_DATA["middleware_competitor"])
                    
                    print("    ✓ Filled financial data and product footprint")
    
    # Save the filled document
    print(f"\nSaving to: {output_path}")
    prs.save(output_path)
    print("✓ Done!")
    
    return output_path


if __name__ == "__main__":
    # Paths
    template_path = Path(__file__).parent.parent / "Executive Briefing Document (16).pptx"
    output_path = Path(__file__).parent.parent / "EBD_Amazon_FILLED.pptx"
    
    if not template_path.exists():
        print(f"Error: Template not found at {template_path}")
        exit(1)
    
    seed_ebd_document(str(template_path), str(output_path))
    
    print("\n" + "="*60)
    print("Now test extraction on the filled document:")
    print(f"  python tools/extract_ebd.py \"{output_path}\"")

