"""
EBD (Executive Briefing Document) PPTX Extractor

Extracts text and table content from EBD PowerPoint files.
Run this standalone to test extraction on a PPTX file.

Usage:
    python tools/extract_ebd.py path/to/ebd.pptx
"""
import json
import sys
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation


def extract_pptx_content(pptx_path: str) -> Dict[str, Any]:
    """
    Extract all text and table content from a PowerPoint file.
    
    Args:
        pptx_path: Path to the PPTX file
        
    Returns:
        Dict containing:
        - slides: List of slide content with text and tables
        - all_text: Flattened list of all text found
        - tables: List of all tables found
    """
    prs = Presentation(pptx_path)
    
    result = {
        "file": str(pptx_path),
        "slide_count": len(prs.slides),
        "slides": [],
        "all_text": [],
        "tables": [],
    }
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_data = {
            "slide_number": slide_num,
            "texts": [],
            "tables": [],
        }
        
        for shape in slide.shapes:
            # Extract text from text frames
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_data["texts"].append(text)
                        result["all_text"].append(text)
            
            # Extract table content
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        row_data.append(cell_text)
                    table_data.append(row_data)
                
                table_entry = {
                    "slide": slide_num,
                    "rows": len(shape.table.rows),
                    "cols": len(shape.table.columns),
                    "data": table_data,
                }
                slide_data["tables"].append(table_entry)
                result["tables"].append(table_entry)
        
        result["slides"].append(slide_data)
    
    return result


def format_extracted_content(extracted: Dict[str, Any]) -> str:
    """
    Format extracted content as readable text for LLM consumption.
    
    Args:
        extracted: Output from extract_pptx_content
        
    Returns:
        Formatted string representation of the EBD content
    """
    lines = []
    lines.append(f"=== EBD CONTENT ({extracted['slide_count']} slides) ===\n")
    
    for slide in extracted["slides"]:
        lines.append(f"\n--- Slide {slide['slide_number']} ---")
        
        # Add texts
        for text in slide["texts"]:
            lines.append(text)
        
        # Add tables
        for table in slide["tables"]:
            lines.append(f"\n[Table {table['rows']}x{table['cols']}]")
            for row in table["data"]:
                # Filter empty cells and join
                non_empty = [c for c in row if c]
                if non_empty:
                    lines.append(" | ".join(non_empty))
    
    return "\n".join(lines)


def extract_ebd_as_text(pptx_path: str) -> str:
    """
    Convenience function: Extract PPTX and return formatted text.
    
    Args:
        pptx_path: Path to the PPTX file
        
    Returns:
        Formatted text representation of EBD content
    """
    extracted = extract_pptx_content(pptx_path)
    return format_extracted_content(extracted)


# CLI for testing
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python tools/extract_ebd.py <path_to_pptx>")
        print("\nExample:")
        print("  python tools/extract_ebd.py ~/Documents/EBD_Acme_Corp.pptx")
        sys.exit(1)
    
    pptx_path = sys.argv[1]
    
    if not Path(pptx_path).exists():
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)
    
    print(f"Extracting content from: {pptx_path}\n")
    
    # Extract raw data
    extracted = extract_pptx_content(pptx_path)
    
    # Print summary
    print(f"Found {extracted['slide_count']} slides")
    print(f"Found {len(extracted['all_text'])} text elements")
    print(f"Found {len(extracted['tables'])} tables")
    print("\n" + "="*60 + "\n")
    
    # Print formatted content
    formatted = format_extracted_content(extracted)
    print(formatted)
    
    # Optionally save JSON output
    if len(sys.argv) > 2 and sys.argv[2] == "--json":
        json_path = Path(pptx_path).stem + "_extracted.json"
        with open(json_path, "w") as f:
            json.dump(extracted, f, indent=2)
        print(f"\n\nJSON saved to: {json_path}")

