"""
EBD Seeding Script for Apple - Testing purposes

Creates a filled-out EBD document tailored to Apple's profile in the database.
Run: python tools/seed_ebd_apple.py
"""
from pptx import Presentation
from pathlib import Path


# Apple-specific dummy data matching their DB profile
APPLE_DATA = {
    "company_name": "Apple",
    "industry": "Not For Profit",  # As per DB
    "address": "One Apple Park Way, Cupertino, CA 95014",
    "website": "www.apple.com",
    
    "meeting_date": "January 20, 2026 | 10:00 AM - 5:00 PM PST",
    "meeting_location": "Oracle Executive Briefing Center - Redwood City",
    "engagement_type": "1:1 Executive Briefing",
    
    "ebd_contact": "Lisa Park, Sr. Account Executive, lisa.park@oracle.com, (650) 555-1234",
    
    "oracle_attendees": """Thomas Reed, VP Analytics Solutions, thomas.reed@oracle.com
Sarah Kim, Sr. Director Customer Experience, sarah.kim@oracle.com
James Liu, Principal Data Architect, james.liu@oracle.com""",
    
    "customer_attendees": """Tim Harrison, Chief Digital Officer
Michelle Torres, VP Marketing Operations
David Chen, CIO
Patricia Wong, Director of Analytics
Ryan Mitchell, Head of Customer Insights""",
    
    "company_description": """Apple is a global technology leader known for innovation in consumer electronics, software, and services. With a focus on customer experience and data-driven decision making, Apple is exploring Oracle Analytics to enhance their marketing operations, customer insights, and service automation capabilities across their global operations.""",
    
    "meeting_objectives": """1. Evaluate Oracle Analytics Cloud for unified customer insights across all touchpoints
2. Understand how Oracle can support Revenue Transformation through data-driven marketing
3. Explore Marketing and Sales Unification capabilities to improve campaign effectiveness
4. Assess Service Automation features for enhancing customer support operations
5. Review integration options with existing Apple data infrastructure""",
    
    "business_challenges": """1. DATA SILOS: Customer data fragmented across 15+ systems including retail, online, support, and services - no unified view of customer journey

2. MARKETING ATTRIBUTION: Difficulty measuring true ROI of marketing campaigns across channels; estimated $50M+ in inefficient ad spend annually

3. SERVICE SCALABILITY: Customer support volume growing 30% YoY; need AI-driven automation to maintain quality while scaling

4. REAL-TIME INSIGHTS: 48-hour lag in getting actionable analytics; competitors making decisions in real-time""",
    
    "it_strategy": """CDO Vision: "Transform Apple into the world's most data-intelligent company by 2027"

IT Strategy:
- Unified data platform for all customer touchpoints
- AI/ML-first approach to customer insights
- Real-time analytics for marketing and service decisions
- Privacy-preserving personalization at scale
- Cloud-native architecture for global scalability""",
    
    "customer_lifecycle": """SALES STAGE: Discovery / Qualified Opportunity ($12M TCV)
- Initial exploration started Q4 2025
- Competitive evaluation with Salesforce Tableau and Microsoft Power BI
- Technical POC planned for Q1 2026
- Decision expected Q2 2026

WHY CONSIDERING ORACLE:
- Impressed by Oracle's autonomous database capabilities
- Strong references from retail industry peers
- Integration capabilities with existing Oracle infrastructure

RECENT IMPLEMENTATIONS:
- Oracle Autonomous Data Warehouse POC (positive results)
- Oracle Cloud Infrastructure evaluation (ongoing)""",
    
    "account_status": """YELLOW - Opportunity with caution

POSITIVES:
- CIO David Chen is Oracle advocate from previous role
- Strong technical interest from analytics team
- Budget approved for 2026 analytics transformation

POTENTIAL DERAILERS:
- CDO has existing relationship with Salesforce executives
- Concerns about Oracle's consumer/retail analytics depth
- Internal preference for "best of breed" vs. platform approach
- Privacy requirements extremely stringent""",
    
    "oracle_talking_points": """1. UNIFIED CUSTOMER 360: Showcase Oracle Analytics' ability to create unified customer views across all touchpoints - reference similar implementations at major retailers achieving 40% improvement in customer insights

2. REVENUE TRANSFORMATION: Demonstrate how Oracle drives marketing ROI through AI-powered attribution and campaign optimization - highlight $30M+ annual savings at comparable companies

3. SERVICE AUTOMATION: Present Oracle's AI-driven service automation capabilities that maintain quality while scaling - emphasize privacy-preserving personalization features critical for Apple""",
    
    # Attendee profiles
    "attendee_1_name": "Tim Harrison, Chief Digital Officer, tim.harrison@apple.com",
    "attendee_1_perspective": "Driving Apple's data transformation agenda. Skeptical of large platform vendors - prefers best-of-breed. Strong relationship with Salesforce. Needs to see clear differentiation.",
    "attendee_1_bio": "20+ years in digital transformation, former CDO at Nike. Harvard MBA. Known for data-driven marketing innovation.",
    
    "attendee_2_name": "David Chen, CIO, david.chen@apple.com",
    "attendee_2_perspective": "Oracle advocate from previous company. Focused on integration and total cost of ownership. Concerned about implementation timeline and change management.",
    "attendee_2_bio": "Former VP of IT at Cisco. 25 years in enterprise technology. Led successful Oracle implementation at previous employer.",
    
    "attendee_3_name": "Michelle Torres, VP Marketing Operations, michelle.torres@apple.com",
    "attendee_3_perspective": "Primary business stakeholder. Frustrated with current analytics delays. Wants real-time campaign optimization and better attribution. Very hands-on, will want to see demos.",
    "attendee_3_bio": "15 years in marketing operations. Former Director at Google. Expert in marketing automation and analytics.",
    
    # Customer references
    "ref_1_name": "Nike",
    "ref_1_industry": "Retail / Consumer",
    "ref_1_website": "www.nike.com",
    "ref_1_products": "Oracle Analytics Cloud, Oracle CX Cloud",
    "ref_1_summary": "Unified customer data across 200+ retail locations and digital channels. Achieved 35% improvement in marketing attribution accuracy and $40M annual savings in ad spend optimization.",
    
    "ref_2_name": "Starbucks",
    "ref_2_industry": "Retail / Food Service",
    "ref_2_website": "www.starbucks.com",
    "ref_2_products": "Oracle Analytics Cloud, Oracle Autonomous Database",
    "ref_2_summary": "Implemented real-time customer insights platform. Reduced analytics latency from 24 hours to 15 minutes. Personalization engine drives 20% of mobile app revenue.",
    
    "ref_3_name": "Target",
    "ref_3_industry": "Retail",
    "ref_3_website": "www.target.com",
    "ref_3_products": "Oracle Analytics Cloud, Oracle Marketing Cloud",
    "ref_3_summary": "Marketing and Sales Unification across digital and physical channels. 45% improvement in campaign effectiveness and real-time inventory-aware promotions.",
    
    # Account team
    "scd_lam": "Lisa Park, Sr. Account Executive, lisa.park@oracle.com, (650) 555-1234",
    "other_am": "Michael Brown, Strategic Account Manager, michael.brown@oracle.com, (650) 555-5678",
    "cse_csm": "Jennifer Lee, Customer Success Manager, jennifer.lee@oracle.com, (650) 555-9012",
    "exec_sponsor": "Robert Martinez, SVP West Region Sales, robert.martinez@oracle.com",
    "impl_partner": "Accenture - Analytics Practice (Primary), Deloitte (Data Integration)",
    
    # LMS/SIA
    "last_audit_date": "June 2024",
    "audit_active": "No",
    "lms_legal": "No",
    "ula_cert": "N/A",
    "lms_contact": "N/A",
    "sia_engaged": "Yes",
    "sia_contact": "Karen White, SIA Manager",
    
    # Financial
    "total_deals_size": "$12M",
    "customer_segment": "Strategic",
    "tam": "$80M",
    "avg_oracle_spend": "$5.2M",
    "annual_revenue": "$383B",
    "share_of_wallet": "0.01%",
    
    # Opportunities
    "opp_cloud_saas": "$6.5M",
    "opp_cloud_saas_prob": "55%",
    "opp_cloud_paas": "$3.5M",
    "opp_cloud_paas_prob": "45%",
    "opp_services": "$2.0M",
    "opp_services_prob": "60%",
    
    # Product footprint
    "apps_support_spend": "$2.1M",
    "apps_footprint": "None",
    "apps_competitor": "Salesforce Marketing Cloud",
    "db_support_spend": "$2.8M",
    "db_footprint": "Oracle Autonomous DW (POC)",
    "db_competitor": "Snowflake, BigQuery",
    "middleware_support_spend": "$0.3M",
    "middleware_footprint": "Minimal",
    "middleware_competitor": "AWS Services",
}


def fill_table_cell(table, row_idx, col_idx, text):
    """Safely fill a table cell if it exists."""
    try:
        if row_idx < len(table.rows) and col_idx < len(table.rows[row_idx].cells):
            cell = table.rows[row_idx].cells[col_idx]
            cell.text = text
            return True
    except Exception as e:
        print(f"Warning: Could not fill cell [{row_idx}][{col_idx}]: {e}")
    return False


def seed_apple_ebd(template_path: str, output_path: str):
    """Fill the EBD template with Apple data."""
    print(f"Loading template: {template_path}")
    prs = Presentation(template_path)
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        
        if slide_num == 1:
            continue  # Skip instructions
        
        tables = [shape.table for shape in slide.shapes if shape.has_table]
        
        # Slide 2: Main EBD
        if slide_num == 2 and tables:
            t = tables[0]
            fill_table_cell(t, 0, 0, APPLE_DATA["company_name"])
            fill_table_cell(t, 0, 2, APPLE_DATA["industry"])
            fill_table_cell(t, 1, 0, APPLE_DATA["address"])
            fill_table_cell(t, 1, 3, APPLE_DATA["meeting_date"])
            fill_table_cell(t, 2, 3, APPLE_DATA["meeting_location"])
            fill_table_cell(t, 3, 3, APPLE_DATA["engagement_type"])
            fill_table_cell(t, 4, 0, APPLE_DATA["website"])
            fill_table_cell(t, 4, 3, APPLE_DATA["ebd_contact"])
            fill_table_cell(t, 5, 3, APPLE_DATA["oracle_attendees"])
            fill_table_cell(t, 6, 3, APPLE_DATA["customer_attendees"])
            fill_table_cell(t, 7, 1, APPLE_DATA["company_description"])
            fill_table_cell(t, 8, 1, APPLE_DATA["meeting_objectives"])
            fill_table_cell(t, 9, 1, APPLE_DATA["business_challenges"])
            fill_table_cell(t, 10, 1, APPLE_DATA["it_strategy"])
            fill_table_cell(t, 11, 1, APPLE_DATA["customer_lifecycle"])
            fill_table_cell(t, 12, 1, APPLE_DATA["account_status"])
            fill_table_cell(t, 13, 1, APPLE_DATA["oracle_talking_points"])
            print("  ✓ Slide 2: Main EBD filled")
        
        # Slide 3: Account team & attendees
        if slide_num == 3 and tables:
            for t in tables:
                rows, cols = len(t.rows), len(t.columns)
                
                if rows == 11 and cols == 5:
                    fill_table_cell(t, 2, 0, "1:1 Executive Briefing")
                    fill_table_cell(t, 2, 1, "November 2025 / Virtual")
                    fill_table_cell(t, 2, 2, "Lisa Park")
                    fill_table_cell(t, 2, 3, "David Chen, Patricia Wong")
                    for row_idx in range(4, rows):
                        cell_text = t.rows[row_idx].cells[0].text.strip().lower()
                        if "scd" in cell_text or "lam" in cell_text:
                            fill_table_cell(t, row_idx, 1, APPLE_DATA["scd_lam"])
                        elif "other account" in cell_text:
                            fill_table_cell(t, row_idx, 1, APPLE_DATA["other_am"])
                        elif "cse" in cell_text or "csm" in cell_text:
                            fill_table_cell(t, row_idx, 1, APPLE_DATA["cse_csm"])
                        elif "executive sponsor" in cell_text:
                            fill_table_cell(t, row_idx, 1, APPLE_DATA["exec_sponsor"])
                        elif "implementation" in cell_text:
                            fill_table_cell(t, row_idx, 1, APPLE_DATA["impl_partner"])
                    print("  ✓ Slide 3: Account team filled")
                
                if rows == 4 and cols == 5:
                    fill_table_cell(t, 1, 0, APPLE_DATA["last_audit_date"])
                    fill_table_cell(t, 1, 1, APPLE_DATA["audit_active"])
                    fill_table_cell(t, 1, 2, APPLE_DATA["lms_legal"])
                    fill_table_cell(t, 1, 3, APPLE_DATA["ula_cert"])
                    fill_table_cell(t, 3, 0, APPLE_DATA["sia_engaged"])
                    fill_table_cell(t, 3, 4, APPLE_DATA["sia_contact"])
                    print("  ✓ Slide 3: LMS/SIA filled")
                
                if rows == 9 and cols == 3:
                    for i, key in enumerate(["attendee_1", "attendee_2", "attendee_3"]):
                        fill_table_cell(t, i*3, 1, APPLE_DATA[f"{key}_name"])
                        fill_table_cell(t, i*3, 2, APPLE_DATA[f"{key}_name"])
                        fill_table_cell(t, i*3+1, 1, APPLE_DATA[f"{key}_perspective"])
                        fill_table_cell(t, i*3+1, 2, APPLE_DATA[f"{key}_perspective"])
                        fill_table_cell(t, i*3+2, 1, APPLE_DATA[f"{key}_bio"])
                        fill_table_cell(t, i*3+2, 2, APPLE_DATA[f"{key}_bio"])
                    print("  ✓ Slide 3: Attendee profiles filled")
        
        # Slide 4: References
        if slide_num == 4 and tables:
            refs = [("ref_1", 0), ("ref_2", 1), ("ref_3", 2)]
            for ref_key, table_idx in refs:
                if table_idx < len(tables):
                    t = tables[table_idx]
                    if len(t.rows) == 6:
                        fill_table_cell(t, 0, 1, APPLE_DATA[f"{ref_key}_name"])
                        fill_table_cell(t, 1, 1, APPLE_DATA[f"{ref_key}_industry"])
                        fill_table_cell(t, 2, 1, APPLE_DATA[f"{ref_key}_website"])
                        fill_table_cell(t, 3, 1, APPLE_DATA[f"{ref_key}_products"])
                        fill_table_cell(t, 5, 1, APPLE_DATA[f"{ref_key}_summary"])
            print("  ✓ Slide 4: References filled")
        
        # Slide 5: Financial
        if slide_num == 5 and tables:
            t = tables[0]
            if len(t.rows) >= 25 and len(t.columns) >= 10:
                fill_table_cell(t, 8, 4, APPLE_DATA["opp_cloud_saas"])
                fill_table_cell(t, 8, 5, APPLE_DATA["opp_cloud_saas_prob"])
                fill_table_cell(t, 9, 4, APPLE_DATA["opp_cloud_paas"])
                fill_table_cell(t, 9, 5, APPLE_DATA["opp_cloud_paas_prob"])
                fill_table_cell(t, 13, 4, APPLE_DATA["opp_services"])
                fill_table_cell(t, 13, 5, APPLE_DATA["opp_services_prob"])
                fill_table_cell(t, 14, 4, APPLE_DATA["total_deals_size"])
                fill_table_cell(t, 14, 9, APPLE_DATA["customer_segment"])
                fill_table_cell(t, 15, 4, APPLE_DATA["tam"])
                fill_table_cell(t, 15, 9, APPLE_DATA["avg_oracle_spend"])
                fill_table_cell(t, 16, 4, APPLE_DATA["annual_revenue"])
                fill_table_cell(t, 16, 9, APPLE_DATA["share_of_wallet"])
                fill_table_cell(t, 19, 3, APPLE_DATA["apps_support_spend"])
                fill_table_cell(t, 19, 5, APPLE_DATA["apps_footprint"])
                fill_table_cell(t, 19, 7, APPLE_DATA["apps_competitor"])
                fill_table_cell(t, 20, 3, APPLE_DATA["db_support_spend"])
                fill_table_cell(t, 20, 5, APPLE_DATA["db_footprint"])
                fill_table_cell(t, 20, 7, APPLE_DATA["db_competitor"])
                print("  ✓ Slide 5: Financial filled")
    
    print(f"\nSaving to: {output_path}")
    prs.save(output_path)
    print("✓ Done!")


if __name__ == "__main__":
    template = Path(__file__).parent.parent / "Executive Briefing Document (16).pptx"
    output = Path(__file__).parent.parent / "EBD_Apple_FILLED.pptx"
    
    if not template.exists():
        print(f"Error: Template not found at {template}")
        exit(1)
    
    seed_apple_ebd(str(template), str(output))

